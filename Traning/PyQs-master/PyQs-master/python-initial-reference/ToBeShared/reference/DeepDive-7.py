.docx
.xlsx
.pptx
PDF
--------------------
#Check popular modules from https://github.com/vinta/awesome-python

###Conversion of documents 

##using comtypes
$ pip install comtypes

#from Word to pdf 
import os
import comtypes.client
import time


#for constant, check https://docs.microsoft.com/en-us/office/vba/api/word.wdsaveformat
#for other API, check https://docs.microsoft.com/en-us/office/vba/api/overview/word
#CreateObject(progid, clsctx=None, machine=None, interface=None, dynamic=False, pServerInfo=None)
word = comtypes.client.CreateObject('Word.Application') #would take time 
word.Visible = True
time.sleep(3)

#should be abs path 
in_file = os.path.abspath('data/example.docx')
out_file1 = os.path.abspath('example.pdf')
out_file2 = os.path.abspath('example2.pdf')

doc=word.Documents.Open(in_file) # open docx file 1
doc.SaveAs(out_file1, FileFormat=17) # conversion, wdFormatPDF = 17
doc.Close() # close docx file 1
word.Visible = False
# convert docx file 2 to pdf file 2
doc = word.Documents.Open(in_file) # open docx file 2
doc.SaveAs(out_file2, FileFormat=17) # conversion, wdFormatPDF = 17
doc.Close() # close docx file 2   
word.Quit() # close Word Application 


#from pptx to pdf 
#for constant, https://docs.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype
#or PowerPoint, start it, press ALT+F11 to open the VBA editor, 
#press F2 to open the Object Browser then search on SaveAs to get this list. 
#Click on any constant name to see the value of the constant at the bottom of the dialog.

#for other methods, check https://docs.microsoft.com/en-us/office/vba/api/overview/powerpoint
ppt = comtypes.client.CreateObject('Powerpoint.Application') #would take time 
ppt.Visible = True
time.sleep(3)
#should be abs path 
in_file = os.path.abspath('data/example.pptx')
out_file1 = os.path.abspath('example.pdf')
out_file2 = os.path.abspath('example2.pdf')

doc=ppt.Presentations.Open(in_file) 
doc.SaveAs(out_file1, FileFormat=32) # conversion, wdFormatPDF = 32
doc.Close() 
#ppt.Visible = False  #not allowed 
doc = ppt.Presentations.Open(in_file) 
doc.SaveAs(out_file2, FileFormat=32) # conversion, wdFormatPDF = 32
doc.Close() 
ppt.Quit() 


#from excel to pdf 
#for constant, check https://docs.microsoft.com/en-us/office/vba/api/excel.xlfileformat
#for other API - check https://docs.microsoft.com/en-us/office/vba/api/overview/excel
xl = comtypes.client.CreateObject('Excel.Application') #would take time 
xl.Visible = True
time.sleep(3)
#should be abs path 
in_file = os.path.abspath('data/empty_book.xlsx')
out_file1 = os.path.abspath('example.pdf')
out_file2 = os.path.abspath('example2.pdf')

books = xl.Workbooks.Open(in_file) 
doc = books.Worksheets[5]  #"charts" #index from 1
doc.SaveAs(out_file1, FileFormat=57) # conversion, wdFormatPDF = 57
doc.Close() 
#xl.Visible = False  #not allowed 
books = xl.Workbooks.Open(in_file) 
doc = books.Worksheets[1]  #"first"
doc.SaveAs(out_file2, FileFormat=57) # conversion, wdFormatPDF = 57
doc.Close() 
xl.Quit() 


#If magic number is not used, you can use enumeration value as well 
from comtypes.client import Constants, CreateObject

powerpoint = CreateObject("Powerpoint.Application")
pp_constants = Constants(powerpoint)
pp_constants.ppSaveAsPDF #32

pres = powerpoint.Presentations.Open(input_path)
pres.SaveAs(output_path, pp_constants.ppSaveAsPDF)
pres.Close() 
powerpoint.Quit() 




###.docx
$ pip install python-docx
#chart functionality, not yet developed,
#convert matplotlib output and embed that 
plt.savefig('some.png')
doc.add_picture('some.png')    
doc.save('myReport.docx')



##Reading Word Documents

import docx
doc = docx.Document('data/example.docx')

>>> dir(doc)
[ 'add_heading', 'add_page_break', 'add_paragraph', 'add_picture',
 'add_section', 'add_table', 'core_properties', 'element', 'inline_shapes', 
 'paragraphs','part', 'save', 'sections', 'settings', 'styles', 'tables']
>>> doc.tables
[]
>>> list(doc.styles)
[_ParagraphStyle('Normal') id: 956857910384, <docx.styles.style._CharacterStyle object at0x000000DEC92D00F0>, ..]
>>> list(doc.paragraphs)
[<docx.text.paragraph.Paragraph object at 0x000000DEC717E7F0>, ..]

>>> dir(docx.text.paragraph.Paragraph)
['add_run', 'alignment', 'clear', 'insert_paragraph_before', 'paragraph_format', 
'part', 'runs', 'style', 'text']

len(doc.paragraphs)
doc.paragraphs[0].text

#Each Paragraph object also has a runs attribute that is a list of Run objects. 
#Run objects have a text attribute and any style info on that text 

len(doc.paragraphs[0].runs)
>>> dir(doc.paragraphs[0].runs[0])
['add_break', 'add_picture', 'add_tab', 'add_text', 'bold', 'clear', 
'element', 'font', 'italic','part', 'style', 'text', 'underline']


##Getting the Full Text from a .docx File

import docx

def getText(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)


print(getText('data/example.docx'))



##Creating docx file 
from docx import Document
from docx.shared import Inches

document = Document()

document.add_heading('Document Title', 0)

p = document.add_paragraph('A plain paragraph having some ')
p.add_run('bold').bold = True
p.add_run(' and some ')
p.add_run('italic.').italic = True
p.add_run('text with emphasis.').style = 'Emphasis' #Applying a character style

document.add_heading('Heading, level 1', level=1)
document.add_paragraph('Intense quote', style='Intense Quote')

document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'first item in ordered list', style='List Number'
)

document.add_picture('data/monty-truth.png', width=Inches(1.25))

records = (
    (3, '101', 'Spam'),
    (7, '422', 'Eggs'),
    (4, '631', 'Spam, spam, eggs, and spam')
)

table = document.add_table(rows=1, cols=3)
dir(table)
#['add_column', 'add_row', 'alignment', 'autofit','cell', 'column_cells', 
#'columns', 'part', 'row_cells', 'rows', 'style', 'table', 'table_direction']

table.style = 'Table Grid'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Qty'
hdr_cells[1].text = 'Id'
hdr_cells[2].text = 'Desc'

#make it bold 
for i in range(3):
    run = hdr_cells[i].paragraphs[0].runs[0]
    run.bold = True
    run.underline = True
    hdr_cells[i].paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

for qty, id, desc in records:
    row_cells = table.add_row().cells
    row_cells[0].text = str(qty)
    row_cells[1].text = id
    row_cells[2].text = desc

document.add_page_break()

document.save('demo.docx')


##Working with Text - Block-level vs. inline text objects

Block Level 
    The paragraph is the primary block-level object in Word.
    A block-level item flows the text it contains between its left and right edges,
    adding an additional line each time the text extends beyond its right boundary.
    For a paragraph, the boundaries are generally the page margins, 
    but they can also be column boundaries if the page is laid out in columns, 
    or cell boundaries if the paragraph occurs inside a table cell.
    A table is also a block-level object.
    The attributes of a block-level item specify its placement on the page, 
    such items as indentation and space before and after a paragraph. 

Inline Text 
    An inline object is a portion of the content that occurs inside a block-level item. 
    An example would be a word that appears in bold or a sentence in all-caps. 
    The most common inline object is a run. 
    All content within a block container is inside of an inline object. 
    Typically, a paragraph contains one or more runs, 
    each of which contain some part of the paragraph’s text.
    The attributes of an inline item generally specify the font in which 
    the content appears, things like typeface, font size, bold, and italic.
    
    
##Paragraph properties
#The formatting properties of a paragraph are accessed using the ParagraphFormat object 
#available using the paragraph's paragraph_format property.

#Horizontal alignment (justification)
from docx.enum.text import WD_ALIGN_PARAGRAPH
document = Document()
paragraph = document.add_paragraph()
paragraph_format = paragraph.paragraph_format

paragraph_format.alignment #None  # indicating alignment is inherited from the style hierarchy
paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
paragraph_format.alignment #CENTER (1)

#Indentation
#Indentation is specified using a Length value, such as Inches, Pt, or Cm. 
#Negative values are valid and cause the paragraph to overlap the margin by the specified amount. 
#Assigning None to an indentation property removes any directly-applied indentation setting 
#and restores inheritance from the style hierarchy:

from docx.shared import Inches
paragraph = document.add_paragraph()
paragraph_format = paragraph.paragraph_format

paragraph_format.left_indent #None  # indicating indentation is inherited from the style hierarchy
paragraph_format.left_indent = Inches(0.5)
paragraph_format.left_indent #457200
paragraph_format.left_indent.inches #0.5

#Right-side indent works in a similar way:

from docx.shared import Pt
paragraph_format.right_indent #None
paragraph_format.right_indent = Pt(24)
paragraph_format.right_indent #304800
paragraph_format.right_indent.pt #24.0

#First-line indent is specified using the first_line_indent property 
#and is interpreted relative to the left indent. 
#A negative value indicates a hanging indent:

paragraph_format.first_line_indent #None
paragraph_format.first_line_indent = Inches(-0.25)
paragraph_format.first_line_indent #-228600
paragraph_format.first_line_indent.inches #-0.25

#Paragraph spacing
paragraph_format.space_before, paragraph_format.space_after #(None, None)  # inherited by default
paragraph_format.space_before = Pt(18)
paragraph_format.space_before.pt #18.0
paragraph_format.space_after = Pt(12)
paragraph_format.space_after.pt #12.0

#Line spacing
#Line spacing is the distance between subsequent baselines in the lines of a paragraph. 
#Line spacing can be specified either as an absolute distance or relative to the line height (essentially the point size of the font used). 
#A typical absolute measure would be 18 points. 
#A typical relative measure would be double-spaced (2.0 line heights). 
#The default line spacing is single-spaced (1.0 line heights).

#line_spacing is either a Length value, a (small-ish) float, or None. 
#A Length value indicates an absolute distance. 
#A float indicates a number of line heights. 
#None indicates line spacing is inherited. 

#line_spacing_rule is a member of the WD_LINE_SPACING enumeration or None:

from docx.shared import Length
paragraph_format.line_spacing #None
paragraph_format.line_spacing_rule #None

paragraph_format.line_spacing = Pt(18)
isinstance(paragraph_format.line_spacing, Length) #True
paragraph_format.line_spacing.pt #18.0
paragraph_format.line_spacing_rule #EXACTLY (4)

paragraph_format.line_spacing = 1.75
paragraph_format.line_spacing #1.75
paragraph_format.line_spacing_rule #MULTIPLE (5)

##Pagination properties

Four paragraph properties, keep_together, keep_with_next, page_break_before, and widow_control 
control aspects of how the paragraph behaves near page boundaries.

keep_together causes the entire paragraph to appear on the same page, 
issuing a page break before the paragraph if it would otherwise be broken across two pages.

keep_with_next keeps a paragraph on the same page as the subsequent paragraph. 
This can be used, for example, to keep a section heading on the same page 
as the first paragraph of the section.

page_break_before causes a paragraph to be placed at the top of a new page. 
This could be used on a chapter heading to ensure chapters start on a new page.

widow_control breaks a page to avoid placing the first or last line of the paragraph 
on a separate page from the rest of the paragraph.

#All four of these properties are tri-state, 
#meaning they can take the value True, False, or None. 
#None indicates the property value is inherited from the style hierarchy. 
#True means “on” and False means “off”:

paragraph_format.keep_together #None  # all four inherit by default
paragraph_format.keep_with_next = True
paragraph_format.keep_with_next #True
paragraph_format.page_break_before = False
paragraph_format.page_break_before #False

##Apply character formatting
#Character formatting is applied at the Run level. 
#Examples include font typeface and size, bold, italic, and underline.


from docx import Document
document = Document()
run = document.add_paragraph().add_run()
font = run.font

#Typeface and size are set like this:
from docx.shared import Pt
font.name = 'Calibri'
font.size = Pt(12)

#Many font properties are tri-state, can take the values True, False, and None. 
#True means the property is 'on', False means it is 'off'. 
#Conceptually, the None value means 'inherit'. 

#A run exists in the style inheritance hierarchy 
#and by default inherits its character formatting from that hierarchy. 
#Any character formatting directly applied using the Font object overrides the inherited values.

#Bold and italic are tri-state properties, 
#as are all-caps, strikethrough, superscript, and many others
font.bold, font.italic #(None, None)
font.italic = True
font.italic #True
font.italic = False
font.italic #False
font.italic = None
font.italic #None

#Underline is a bit of a special case. 
#It is a hybrid of a tri-state property and an enumerated value property. 
#True means single underline. False means no underline, 
#double or dashed are specified with a member of the WD_UNDERLINE enumeration:

font.underline #None
font.underline = True
# or perhaps
font.underline = WD_UNDERLINE.DOT_DASH

#Font color, ColorFormat

from docx.shared import RGBColor
font.color.rgb = RGBColor(0x42, 0x24, 0xE9)

#A font can also be set to a theme color by assigning a member of the MSO_THEME_COLOR_INDEX enumeration:
from docx.enum.dml import MSO_THEME_COLOR
font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

#A font's color can be restored to its default (inherited) value 
font.color.rgb = None

#Determining the color of a font begins with determining its color type:
font.color.type # RGB (1) , a MSO_COLOR_TYPE enumeration

#When the color type is MSO_COLOR_TYPE.RGB
font.color.rgb #RGBColor(0x42, 0x24, 0xe9)

#When the color type is MSO_COLOR_TYPE.THEME
font.color.theme_color #ACCENT_1 (5)

##Working with Sections
#Word supports the notion of a section, 
#a division of a document having the same page layout settings, 
#such as margins and page orientation. 

#This is how, for example, a document can contain some pages in portrait layout and others in landscape.

#Most Word documents have only the single section that comes by default 
document = Document()
sections = document.sections
sections    #<docx.parts.document.Sections object at 0x1deadbeef>
len(sections) #3
section = sections[0]
section #<docx.section.Section object at 0x1deadbeef>
for section in sections:
    print(section.start_type)

#Output 
NEW_PAGE (2)
EVEN_PAGE (3)
ODD_PAGE (4)

#Adding a new section
current_section = document.sections[-1]  # last section in document
current_section.start_type  #NEW_PAGE (2)
new_section = document.add_section(WD_SECTION.ODD_PAGE)
new_section.start_type  #ODD_PAGE (4)

##Section properties

#Section.start_type describes the type of break that precedes the section:
section.start_type  #NEW_PAGE (2)
section.start_type = WD_SECTION.ODD_PAGE
section.start_type  #ODD_PAGE (4)

#Page dimensions and orientation
section.orientation, section.page_width, section.page_height #(PORTRAIT (0), 7772400, 10058400)  # (Inches(8.5), Inches(11))
new_width, new_height = section.page_height, section.page_width
section.orientation = WD_ORIENT.LANDSCAPE
section.page_width = new_width
section.page_height = new_height
section.orientation, section.page_width, section.page_height #(LANDSCAPE (1), 10058400, 7772400)

#Page margins

from docx.shared import Inches
section.left_margin, section.right_margin #(1143000, 1143000)  # (Inches(1.25), Inches(1.25))
section.top_margin, section.bottom_margin   #(914400, 914400)  # (Inches(1), Inches(1))
section.gutter #0
section.header_distance, section.footer_distance #(457200, 457200)  # (Inches(0.5), Inches(0.5))
section.left_margin = Inches(1.5)
section.right_margin = Inches(1)
section.left_margin, section.right_margin   #(1371600, 914400)

##Working with Headers and Footers
#Headers and footers are linked to a section

document = Document()
section = document.sections[0]
header = section.header
header #<docx.section._Header object at 0x...>

#A _Header object is always present on Section.header, 
#even when no header is defined for that section. 

#The presence of an actual header definition is indicated by _Header.is_linked_to_previous:
header.is_linked_to_previous #True
#A value of True indicates the _Header object contains no header definition 
#and the section will display the same header as the previous section. 

#Adding a header definition (general case)
header.is_linked_to_previous #True
header.is_linked_to_previous = False
header.is_linked_to_previous    #False
#The newly added header definition contains a single empty paragraph. 

#Adding a header (simple case)
#Note that like a new document, a new header already contains a single (empty) paragraph:
paragraph = header.paragraphs[0]
paragraph.text = "Title of my document"

#Note also that the act of adding content 
#changed the state of .is_linked_to_previous:
header.is_linked_to_previous #False

#Removing a header
header.is_linked_to_previous = True
header.is_linked_to_previous #True


##Understanding Styles
#Word has paragraph styles, character styles, table styles, and numbering definitions. 
#These are applied to a paragraph, a span of text, a table, and a list, respectively.

style definition
    A <w:style> element in the styles part of a document 
    that explicitly defines the attributes of a style.
defined style
    A style that is explicitly defined in a document. 
    Contrast with latent style.
built-in style
    One of the set of 276 pre-set styles built into Word, such as 'Heading 1'. 
    A built-in style can be either defined or latent. 
    A built-in style that is not yet defined is known as a latent style. 
    Both defined and latent built-in styles may appear as options in Word's style panel and style gallery.
custom style
    Also known as a user defined style, 
    any style defined in a Word document that is not a built-in style. 
    Note that a custom style cannot be a latent style.
latent style
    A built-in style having no definition in a particular document is known 
    as a latent style in that document. 
    A latent style can appear as an option in the Word UI depending on the settings in the LatentStyles object for the document.
recommended style list
    A list of styles that appears in the styles toolbox or panel 
    when 'Recommended' is selected from the 'List:' dropdown box.
Style Gallery
    The selection of example styles that appear in the ribbon of the Word UI 
    and which may be applied by clicking on one of them.

##Identifying a style
#A style has three identifying properties, name, style_id, and type.

#Each style's name property is its stable, unique identifier for access purposes.
#A style's style_id is used internally 
#A style's type is set at creation time and cannot be changed.


#A style can inherit properties from another style, 
#similarly to how Cascading Style Sheets (CSS) works. 
#Inheritance is specified using the base_style attribute.


#If you apply a style using python-docx that’s not defined in your file (in the styles.xml ), Word just ignores it 
#If you apply a style, delete the content you applied it to, 
#and then save the document; the style definition stays in the saved file.

Paragraph styles in default template
    Normal
    Body Text
    Body Text 2
    Body Text 3
    Caption
    Heading 1
    Heading 2
    Heading 3
    Heading 4
    Heading 5
    Heading 6
    Heading 7
    Heading 8
    Heading 9
    Intense Quote
    List
    List 2
    List 3
    List Bullet
    List Bullet 2
    List Bullet 3
    List Continue
    List Continue 2
    List Continue 3
    List Number
    List Number 2
    List Number 3
    List Paragraph
    Macro Text
    No Spacing
    Quote
    Subtitle
    TOCHeading
    Title
Character styles in default template
    Body Text Char
    Body Text 2 Char
    Body Text 3 Char
    Book Title
    Default Paragraph Font
    Emphasis
    Heading 1 Char
    Heading 2 Char
    Heading 3 Char
    Heading 4 Char
    Heading 5 Char
    Heading 6 Char
    Heading 7 Char
    Heading 8 Char
    Heading 9 Char
    Intense Emphasis
    Intense Quote Char
    Intense Reference
    Macro Text Char
    Quote Char
    Strong
    Subtitle Char
    Subtle Emphasis
    Subtle Reference
    Title Char
Table styles in default template
    Table Normal
    Colorful Grid
    Colorful Grid Accent 1
    Colorful Grid Accent 2
    Colorful Grid Accent 3
    Colorful Grid Accent 4
    Colorful Grid Accent 5
    Colorful Grid Accent 6
    Colorful List
    Colorful List Accent 1
    Colorful List Accent 2
    Colorful List Accent 3
    Colorful List Accent 4
    Colorful List Accent 5
    Colorful List Accent 6
    Colorful Shading
    Colorful Shading Accent 1
    Colorful Shading Accent 2
    Colorful Shading Accent 3
    Colorful Shading Accent 4
    Colorful Shading Accent 5
    Colorful Shading Accent 6
    Dark List
    Dark List Accent 1
    Dark List Accent 2
    Dark List Accent 3
    Dark List Accent 4
    Dark List Accent 5
    Dark List Accent 6
    Light Grid
    Light Grid Accent 1
    Light Grid Accent 2
    Light Grid Accent 3
    Light Grid Accent 4
    Light Grid Accent 5
    Light Grid Accent 6
    Light List
    Light List Accent 1
    Light List Accent 2
    Light List Accent 3
    Light List Accent 4
    Light List Accent 5
    Light List Accent 6
    Light Shading
    Light Shading Accent 1
    Light Shading Accent 2
    Light Shading Accent 3
    Light Shading Accent 4
    Light Shading Accent 5
    Light Shading Accent 6
    Medium Grid 1
    Medium Grid 1 Accent 1
    Medium Grid 1 Accent 2
    Medium Grid 1 Accent 3
    Medium Grid 1 Accent 4
    Medium Grid 1 Accent 5
    Medium Grid 1 Accent 6
    Medium Grid 2
    Medium Grid 2 Accent 1
    Medium Grid 2 Accent 2
    Medium Grid 2 Accent 3
    Medium Grid 2 Accent 4
    Medium Grid 2 Accent 5
    Medium Grid 2 Accent 6
    Medium Grid 3
    Medium Grid 3 Accent 1
    Medium Grid 3 Accent 2
    Medium Grid 3 Accent 3
    Medium Grid 3 Accent 4
    Medium Grid 3 Accent 5
    Medium Grid 3 Accent 6
    Medium List 1
    Medium List 1 Accent 1
    Medium List 1 Accent 2
    Medium List 1 Accent 3
    Medium List 1 Accent 4
    Medium List 1 Accent 5
    Medium List 1 Accent 6
    Medium List 2
    Medium List 2 Accent 1
    Medium List 2 Accent 2
    Medium List 2 Accent 3
    Medium List 2 Accent 4
    Medium List 2 Accent 5
    Medium List 2 Accent 6
    Medium Shading 1
    Medium Shading 1 Accent 1
    Medium Shading 1 Accent 2
    Medium Shading 1 Accent 3
    Medium Shading 1 Accent 4
    Medium Shading 1 Accent 5
    Medium Shading 1 Accent 6
    Medium Shading 2
    Medium Shading 2 Accent 1
    Medium Shading 2 Accent 2
    Medium Shading 2 Accent 3
    Medium Shading 2 Accent 4
    Medium Shading 2 Accent 5
    Medium Shading 2 Accent 6
    Table Grid


#Access a style
document = Document()
styles = document.styles
styles #<docx.styles.styles.Styles object at 0x10a7c4f50>

#The Styles object provides dictionary-style access to defined styles by name:
styles['Normal']    #<docx.styles.style._ParagraphStyle object at <0x10a7c4f6b>

#The Styles object is iterable. 
#By using the identification properties on BaseStyle, 
#various subsets of the defined styles can be generated. 

#For example, this code will produce a list of the defined paragraph styles:

from docx.enum.style import WD_STYLE_TYPE
styles = document.styles
paragraph_styles = [
    s for s in styles if s.type == WD_STYLE_TYPE.PARAGRAPH
]
for style in paragraph_styles:
    print(style.name)
    
#Output 
Normal
Body Text
List Bullet

#Apply a style
#The Paragraph, Run, and Table objects each have a style attribute. 

document = Document()
paragraph = document.add_paragraph()
paragraph.style #<docx.styles.style._ParagraphStyle object at <0x11a7c4c50>
paragraph.style.name    #'Normal'
paragraph.style = document.styles['Heading 1']
paragraph.style.name    #'Heading 1'

#A style name can also be assigned directly
paragraph.style = 'List Bullet'
paragraph.style #<docx.styles.style._ParagraphStyle object at <0x10a7c4f84>
paragraph.style.name    #'List Bullet'

#A style can also be applied at creation time 
paragraph = document.add_paragraph(style='Body Text')
paragraph.style.name    #'Body Text'
body_text_style = document.styles['Body Text']
paragraph = document.add_paragraph(style=body_text_style)
paragraph.style.name    #'Body Text'

#Add or delete a style
#A new style can be added to the document by specifying a unique name and a style type:

from docx.enum.style import WD_STYLE_TYPE
styles = document.styles
style = styles.add_style('Citation', WD_STYLE_TYPE.PARAGRAPH)
style.name  #'Citation'
style.type  #PARAGRAPH (1)

#Use the base_style property to specify a style the new style should inherit formatting settings from:
style.base_style    #None
style.base_style = styles['Normal']
style.base_style    #<docx.styles.style._ParagraphStyle object at 0x10a7a9550>
style.base_style.name   #'Normal'

#A style can be removed from the document by calling its delete() method:
#It does not affect content in the document to which that style is applied
#Content having a style not defined in the document is rendered using the default style for that content object, 
#e.g. 'Normal' in the case of a paragraph.
styles = document.styles
len(styles) #10
styles['Citation'].delete()
len(styles) #9

#Define character formatting
from docx import Document
document = Document()
style = document.styles['Normal']
font = style.font

#Typeface and size are set like this:

from docx.shared import Pt
font.name = 'Calibri'
font.size = Pt(12)

#Many font properties are tri-state
font.bold, font.italic  #(None, None)
font.italic = True
font.italic #True
font.italic = False
font.italic #False
font.italic = None
font.italic #None

#Underline is a bit of a special case
font.underline  #None
font.underline = True
# or perhaps
font.underline = WD_UNDERLINE.DOT_DASH


#Define paragraph formatting
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Inches, Pt
document = Document()
style = document.styles.add_style('Indent', WD_STYLE_TYPE.PARAGRAPH)
paragraph_format = style.paragraph_format
paragraph_format.left_indent = Inches(0.25)
paragraph_format.first_line_indent = Inches(-0.25)
paragraph_format.space_before = Pt(12)
paragraph_format.widow_control = True

#Use paragraph-specific style properties
#A paragraph style has a next_paragraph_style property 
#that specifies the style to be applied to new paragraphs inserted 
#after a paragraph of that style. 
from docx import Document
document = Document()
styles = document.styles

styles['Heading 1'].next_paragraph_style = styles['Body Text']

#The default behavior can be restored by assigning None or the style itself:
heading_1_style = styles['Heading 1']
heading_1_style.next_paragraph_style.name   #'Body Text'

heading_1_style.next_paragraph_style = heading_1_style
heading_1_style.next_paragraph_style.name   #'Heading 1'

heading_1_style.next_paragraph_style = None
heading_1_style.next_paragraph_style.name   #'Heading 1'










###.pptx
$ pip install python-pptx
#You can open any PowerPoint 2007 or later file


##Extract all text from slides in presentation
#Presentation Consists of slides , each slide might created from layout(created in master slide)
#layout is composed of many preformatted placeholders(kind of shape)
#OR a slide can be created from many Shapes 
#Each slide has a shape tree(.shapes) that holds its shapes(created directly or via layout)


#Presentation has list of slides, each slide has list of shapes 
#Each shape might have text_frame, text_frame has list of paragraphs
#Each paragraph has list of runs, each run has .text 

from pptx import Presentation

prs = Presentation('data/example.pptx')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

##Hello World example

from pptx import Presentation

prs = Presentation()

>>> list(prs.slide_layouts)
[<pptx.slide.SlideLayout object at 0x0000003E2451EC28>, ...]
>>> len(prs.slide_layouts)
11
>>> dir(pptx.slide.SlideLayout)
['background', 'element', 'iter_cloneable_placeholders', 'name', 'part', 
'placeholders', 'shapes', 'slide_master', 'used_by_slides']
>>> { i: s.name for i,s in enumerate(prs.slide_layouts)}
{0: 'Title Slide', 1: 'Title and Content', 2: 'Section Header', 
3: 'Two Content', 4: 'Comparison', 5: 'Title Only', 6: 'Blank', 
7: 'Content with Caption', 8: 'Picture with Caption', 
9: 'Title and Vertical Text', 10: 'Vertical Title and Text'}

>>> dir(slide.shapes)
['add_chart', 'add_connector', 'add_group_shape', 'add_movie', 'add_picture', 
'add_shape', 'add_table', 'add_textbox', 'build_freeform', 
'clone_layout_placeholders', 'clone_placeholder', 'element', 'index', 'parent', 
'part', 'ph_basename', 'placeholders', 'title', 'turbo_add_enabled']

#Add title and some text 
prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "first automated pptx!"

#Adding a bullet slide 
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title

body = slide.shapes.placeholders[1]

title.text = 'Adding a Bullet Slide'

tf = body.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use text_frame.text for above bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use text_frame.add_paragraph() for subsequent bullets'
p.level = 2

#add textbox

from pptx.util import Inches, Pt

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
body = slide.shapes.add_textbox(left, top, width, height)
tf = body.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

#add_picture() example

img_path = 'data/monty-truth.png'

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

#add_shape() example
#Constants representing each of the available auto shapes (like MSO_SHAPE.ROUNDED_RECT, MSO_SHAPE.CHEVRON, etc.) are listed on the autoshape-types page.
from pptx.enum.shapes import MSO_SHAPE

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

slide.shapes.title.text = 'Adding an AutoShape'

left = Inches(0.93)  # 0.93" centers this overall set of shapes
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

body = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
body.text = 'Step 1'

left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    body = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    body.text = 'Step %d' % n
    left = left + width - Inches(0.4)

#add_table() example

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

slide.shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

#make it bold and Center
from pptx.enum.text import PP_ALIGN
for i in range(2):
    run = table.cell(0, i).text_frame.paragraphs[0].runs[0]
    run.bold = True
    run.underline = True
    table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    
# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

#Add a chart 
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Adding a chart'

chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

from pptx.enum.chart import XL_TICK_MARK

category_axis = chart.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)

value_axis = chart.value_axis
value_axis.maximum_scale = 50.0
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"%"'
tick_labels.font.bold = True
tick_labels.font.size = Pt(14)

prs.save('test.pptx')






 
    

##Presentation
#Presentation Consists of slides , each slide might created from layout(created in master slide)
#layout is composed of many preformated placeholders(kind of shape)
#OR a slide can be created from many Shapes 
#Each slide has a shape tree(.shapes) that holds its shapes(created directly or via layout)


# start with default presentation
prs = Presentation()

##Slide masters
#A presentation has a list of slide masters and a list of slides. 
#When you want all your slides to contain the same fonts and images (such as logos), you can make those changes in one place—the Slide Master, 
#and they'll be applied to all your slides. 

from pptx import Presentation

prs = Presentation()
slide_master = prs.slide_masters[0]
# is equivalent to
slide_master = prs.slide_master

##Slide layouts
#Slide layouts belong to a slide master, not directly to a presentation, 
#But , can also be accessed directly from the presentation via syntactic sugar:

prs = Presentation()
title_slide_layout = prs.slide_masters[0].slide_layouts[0]
# is equivalent to:
title_slide_layout = prs.slide_layouts[0]

##Slides
prs = Presentation(path)
first_slide = prs.slides[0]


  
    
    
    
##Slide layout basics
#A slide layout is like a template for a slide. 

#The slide layouts in a standard PowerPoint theme always occur in the same sequence.
#This allows content from one deck to be pasted into another 
#and be connected with the right new slide layout:
    Title (presentation title slide)
    Title and Content
    Section Header (sometimes called Segue)
    Two Content (side by side bullet textboxes)
    Comparison (same but additional title for each side by side content box)
    Title Only
    Blank
    Content with Caption
    Picture with Caption


##Adding a slide
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
new_slide = prs.slides.add_slide(title_slide_layout)


##Understanding Shapes
#Pretty much anything on a slide is a shape;( except slide background)

#Technically there are only six different types of shapes that can be placed on a slide:
auto shape
    This is a regular shape, like a rectangle, an ellipse, or a block arrow. 
    They come in a large variety of preset shapes, 
    in the neighborhood of 180 different ones. 
    An auto shape can have a fill and an outline, and can contain text. 
    Some auto shapes have adjustments, the little yellow diamonds you can drag 
    to adjust how round the corners of a rounded rectangle are for example. 
    A text box is also an autoshape, a rectangular one, 
    just by default without a fill and without an outline.
picture
    A raster image, like a photograph or clip art is referred to 
    as a picture in PowerPoint. 
    Its its own kind of shape with different behaviors than an autoshape. 
    Note that an auto shape can have a picture fill, 
    in which an image 'shows through' as the background of the shape 
    instead of a fill color or gradient. 
    That's a different thing. 
graphic frame
    This is the technical name for the container 
    that holds a table, a chart, a smart art diagram, or media clip. 
    You cant add one of these by itself, it just shows up in the file 
    when you add a graphical object. 
group shape
    In PowerPoint, a set of shapes can be grouped, 
    allowing them to be selected, moved, resized, and even filled as a unit. 
    When you group a set of shapes a group shape gets created 
    to contain those member shapes. 
    You cant actually see these except by their bounding box 
    when the group is selected.
line/connector
    Lines are different from auto shapes because, they are linear. 
    Some lines can be connected to other shapes and stay connected 
    when the other shape is moved. 
    These are not supported yet 
content part
    It has something to do with embedding 'foreign' XML like SVG 
    in with the presentation.
    
    
#Practically, we use below types 
shape shapes – auto shapes with fill and an outline
text boxes – auto shapes with no fill and no outline
placeholders – auto shapes that can appear on a slide layout or master 
               and be inherited on slides that use that layout, 
               allowing content to be added that takes on the formatting 
               of the placeholder
line/connector – as described above
picture – as described above
table – that row and column thing
chart – pie chart, line chart, etc.
smart art – not supported yet, although preserved if present
media clip – video or audio

##Accessing the shapes on a slide
#Each slide has a shape tree that holds its shapes. 
shapes = slide.shapes


##Working with AutoShapes
#Auto shapes are regular shape shapes. Squares, circles, triangles, stars, ..
#There are 182 different auto shapes to choose from. 
#120 of these have adjustment 'handles' 
#https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html

##Adding an auto shape
#The following code adds a rounded rectangle shape,
#one inch square, and positioned one inch from the top-left corner of the slide:

from pptx.enum.shapes import MSO_SHAPE

shapes = slide.shapes
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(  
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

##Understanding English Metric Units
#Internally, PowerPoint stores length values in English Metric Units (EMU). 
#EMU is an integer unit of length, 914400 to the inch. 
#914400 has the great virtue that it is evenly divisible by a great many common factors, 
#allowing exact conversion between inches and centimeters, 
#for example. Being an integer, it can be represented exactly across serializations and across platforms.


from pptx.util import Inches, Pt
length = Inches(1)
length #914400
length.inches #1.0
length.cm #2.54
length.pt #72.0
length = Pt(72)
length #914400

##Shape position and dimensions
#All shapes have a position on their slide and have a size. 

from pptx.enum.shapes import MSO_SHAPE
left = top = width = height = Inches(1.0)
shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)
shape.left, shape.top, shape.width, shape.height #(914400, 914400, 914400, 914400)
shape.left.inches #1.0
shape.left = Inches(2.0) #shape.left.inches
2.0

##Fill
#AutoShapes have an outline around their outside edge. 
#What appears within that outline is called the shape's fill.

#The most common type of fill is a solid color. 
#A shape may also be filled with a gradient, a picture, a pattern (like cross-hatching for example), or may have no fill (transparent).

#This code sets the fill of a shape to red:
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

#This sets it to the theme color that appears as 'Accent 1 - 25% Darker' 
from pptx.enum.dml import MSO_THEME_COLOR
fill = shape.fill
fill.solid()
fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
fill.fore_color.brightness = -0.25

#This sets the shape fill to transparent, 
#or 'No Fill' as it's called in the PowerPoint UI:
shape.fill.background()

#the first step is to specify the desired fill type 
#by calling the corresponding method on fill. 
#Doing so actually changes the properties available on the fill object
fill = shape.fill
fill.solid()
fill.fore_color #<pptx.dml.color.ColorFormat object at 0x10ce20910>
fill.background()
fill.fore_color
Traceback (most recent call last):
  ...
TypeError: a transparent (background) fill has no foreground color

##Line
#The outline of an AutoShape can also be formatted, 
#including setting its color, width, dash (solid, dashed, dotted, etc.), 
#line style (single, double, thick-thin, etc.), end cap, join type, and others. 

#color and width can be set using python-pptx:

line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
line.color.brightness = 0.5  # 50% lighter
line.width = Pt(2.5)

#Theme colors can be used on lines too:
line.color.theme_color = MSO_THEME_COLOR.ACCENT_6

#Shape.line has the attribute .color. 
#This is essentially a shortcut for:

line.fill.solid()
line.fill.fore_color

#This makes sense for line formatting 
#because a shape outline is most frequently set to a solid color. 
#Accessing the fill directly is required, 
#for example, to set the line to transparent:
line.fill.background()

##Line width
#The shape outline also has a read/write width property:

line.width  #9525
line.width.pt   #0.75
line.width = Pt(2.0)
line.width.pt   #2.0

##Shape Adjustment Concepts

#adjustments are particular to a specific auto shape type. 
#Each auto shape has between zero and eight adjustments. 
#What each of them does is arbitrary and depends on the shape design.

#Adjustment values are large integers,
# each based on a nominal value of 100,000. 
#The effective value of an adjustment is proportional to the width or height of the shape. 

#So a value of 50,000 for an x-coordinate adjustment corresponds 
#to half the width of the shape; 
#a value of 75,000 for a y-coordinate adjustment corresponds to 3/4 of the shape height.

#Adjustment values can be negative, generally indicating the coordinate is 
#to the left or above the top left corner (origin) of the shape. 

#Values can also be subject to limits, 
#meaning their effective value cannot be outside a prescribed range. 
#In practice this corresponds to a point not being able to extend beyond the left side of the shape, for example.

#The following code formats a callout shape using its adjustments:

callout_sp = shapes.add_shape(
    MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR, left, top, width, height
)

# get the callout line coming out of the right place
adjs = callout_sp.adjustments
adjs[0] = 0.5   # vert pos of junction in margin line, 0 is top
adjs[1] = 0.0   # horz pos of margin ln wrt shape width, 0 is left side
adjs[2] = 0.5   # vert pos of elbow wrt margin line, 0 is top
adjs[3] = -0.1  # horz pos of elbow wrt shape width, 0 is margin line
adjs[4] = 3.0   # vert pos of line end wrt shape height, 0 is top
a5 = adjs[3] - (adjs[4] - adjs[0]) * height/width
adjs[5] = a5    # horz pos of elbow wrt shape width, 0 is margin line

# rotate 45 degrees counter-clockwise
callout_sp.rotation = -45.0


##Understanding placeholders
#a placeholder is a pre-formatted shape into which content can be placed. 

#the auto shape (p:sp element), picture (p:pic element), and graphic frame (p:graphicFrame) shape types can be a placeholder. 
#The group shape (p:grpSp), connector (p:cxnSp), and content part (p:contentPart) shapes cannot be a placeholder. 
#A graphic frame placeholder can contain a table, a chart, or SmartArt.

##Placeholder types
#There are 18 types of placeholder.

Title, Center Title, Subtitle, Body
    These placeholders typically appear on a conventional 'word chart' 
    containing text only, often organized as a title 
    and a series of bullet points. 
    All of these placeholders can accept text only.
Content
    This multi-purpose placeholder is the most commonly used 
    for the body of a slide. 
    When unpopulated, it displays 6 buttons to allow insertion of 
    a table, a chart, SmartArt, a picture, clip art, or a media clip.
Picture, Clip Art
    These both allow insertion of an image. 
    The insert button on a clip art placeholder brings up the clip art gallery 
    rather than an image file chooser, but otherwise these behave the same.
Chart, Table, Smart Art
    These three allow the respective type of rich graphical content 
    to be inserted.
Media Clip
    Allows a video or sound recording to be inserted.
Date, Footer, Slide Number
    These three appear on most slide masters and slide layouts, 
    but do not behave as most users would expect. 
    These also commonly appear on the Notes Master and Handout Master.
Header
    Only valid on the Notes Master and Handout Master.
Vertical Body, Vertical Object, Vertical Title
    Used with vertically oriented languages such as Japanese.


##Placholders inherit
#A placeholder appearing on a slide is only part of the overall placeholder mechanism. 
#Placeholder behavior requires three different categories of placeholder shape; 
#those that exist on a slide master, those on a slide layout, 
#and those that ultimately appear on a slide in a presentation.

#These three categories of placeholder participate in a property inheritance hierarchy, 
#either as an inheritor, an inheritee, or both. 

#Placeholder shapes on masters are inheritees only. 
#Conversely placeholder shapes on slides are inheritors only. 
#Placeholders on slide layouts are both, a possible inheritor from a slide master placeholder 
#and an inheritee to placeholders on slides linked to that layout.

#A layout inherits from its master differently than a slide inherits from its layout. 
#A layout placeholder inherits from the master placeholder sharing the same type. 
#A slide placeholder inherits from the layout placeholder having the same idx value.

#In general, all formatting properties are inherited from the 'parent' placeholder. 
#This includes position and size as well as fill, line, and font. 
#Any directly applied formatting overrides the corresponding inherited value. 
#Directly applied formatting can be removed be reapplying the layout.


##Working with placeholders
#Placeholders can make adding content a lot easier. 
#If you've ever added a new textbox to a slide from scratch 
#and noticed how many adjustments it took to get it the way you wanted 

#The placeholder is in the right position with the right font size, 
#paragraph alignment, bullet style, etc., etc. 
#Basically you can just click and type in some text and you've got a slide.

#A placeholder can be also be used to place a rich-content object on a slide. 
#A picture, table, or chart can each be inserted into a placeholder 
#and so take on the position and size of the placeholder, 
#as well as certain of its formatting attributes.

##Access a placeholder

#The most reliable way to access a known placeholder is by its idx value. 
#The idx value of a placeholder is the integer key of the slide layout placeholder it inherits properties from. 

#As such, it remains stable throughout the life of the slide 
#and will be the same for any slide created using that layout.

#In general, the idx value of a placeholder from a built-in slide layout 
#will be between 0 and 5. 
#The title placeholder will always have idx 0 if present 
#and any other placeholders will follow in sequence, top to bottom and left to right. 
#A placeholder added to a slide layout by a user in PowerPoint will receive an idx value starting at 10.


#Take a look at the placeholders on a slide and pick out the one you want:

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

#Output 
0  Title 1
1  Picture Placeholder 2
2  Text Placeholder 3

#then, having the known index in hand, to access it directly
#Note it is a dict with integer as key 
slide.placeholders[1]   #<pptx.parts.slide.PicturePlaceholder object at 0x10d094590>
slide.placeholders[2].name  #'Text Placeholder 3'

##Identify and Characterize a placeholder
#A placeholder behaves differently that other shapes in some ways. 
#In particular, the value of its shape_type attribute is unconditionally MSO_SHAPE_TYPE.PLACEHOLDER 

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
for shape in slide.shapes:
    print('%s' % shape.shape_type)

#Output 
PLACEHOLDER (14)
PLACEHOLDER (14)
PLACEHOLDER (14)

#To find out more
for shape in slide.shapes:
    if shape.is_placeholder:
        phf = shape.placeholder_format
        print('%d, %s' % (phf.idx, phf.type))

#output 
0, TITLE (1)
1, PICTURE (18)
2, BODY (2)

#Another way a placeholder acts differently is that it inherits 
#its position and size from its layout placeholder. 
#This inheritance is overridden if the position and size of a placeholder are changed.


##Insert content into a placeholder
#Certain placeholder types have specialized methods for inserting content. 

#the picture, table, and chart placeholders have content insert_* methods. 
#Text can be inserted into title 
#and body placeholders in the same way text is inserted into an auto shape.

#The picture placeholder has an insert_picture() method:

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
placeholder = slide.placeholders[1]  # idx key, not position
placeholder.name  #'Picture Placeholder 2'
placeholder.placeholder_format.type #PICTURE (18)
#A reference to a picture 'placeholder' becomes invalid after its insert_picture() method is called. 
picture = placeholder.insert_picture('data/monty-truth.png')
#A picture inserted in this way is stretched proportionately 
#and cropped to fill the entire placeholder. 


##TablePlaceholder.insert_table()
#The table placeholder has an insert_table() method. 

#The built-in template has no layout containing a table placeholder, 
#so create one 
prs = Presentation('data/having-table-placeholder.pptx')

for e in prs.slide_layouts:
    print(e.name)
    for e1 in e.placeholders:
            print("\t",e1.name, e1.placeholder_format.idx)

#Example 
slide = prs.slides.add_slide(prs.slide_layouts[11])
placeholder = slide.placeholders[13]  # idx key, not position
placeholder.name  #'Table Placeholder 1'
placeholder.placeholder_format.type  #TABLE (12)
#a reference to a table placeholder becomes invalid after its insert_table() method is called
graphic_frame = placeholder.insert_table(rows=2, cols=2)
table = graphic_frame.table
len(table.rows), len(table.columns) #(2, 2)

table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


##ChartPlaceholder.insert_chart()
#The chart placeholder has an insert_chart() method.
 
#The presentation template built into python-pptx has no layout containing a chart placeholder,
#so create one 
prs = Presentation('data/having-chart-placeholder.pptx')

for e in prs.slide_layouts:
    print(e.name)
    for e1 in e.placeholders:
            print("\t",e1.name, e1.placeholder_format.idx)
#code 
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation('data/having-chart-placeholder.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[0])

placeholder = slide.placeholders[13]  # idx key, not position
placeholder.name    #'Chart Placeholder 9'
placeholder.placeholder_format.type #CHART (12)

chart_data = ChartData()
chart_data.categories = ['Yes', 'No']
chart_data.add_series('Series 1', (42, 24))
#a reference to a chart placeholder becomes invalid after its insert_chart() method is called
graphic_frame = placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data)
chart = graphic_frame.chart
chart.chart_type #PIE (5)


##Setting the slide title
#Almost all slide layouts have a title placeholder, 
#is a common operation and there's a dedicated attribute on the shape tree for it:
title_placeholder = slide.shapes.title
title_placeholder.text = 'Air-speed Velocity of Unladen Swallows'


##Working with text
#Auto shapes and table cells can contain text. 
#Other shapes can't. 
#Text is always manipulated the same way, regardless of its container.

#Text exists in a hierarchy of three levels:
    Shape.text_frame
    TextFrame.paragraphs
    _Paragraph.runs

#All the text in a shape is contained in its text frame. 
#A text frame has vertical alignment, margins, wrapping and auto-fit behavior, 
#a rotation angle, some possible 3D visual features, 
#and can be set to format its text into multiple columns. 

#It also contains a sequence of paragraphs, 
#which always contains at least one paragraph, even when empty.
#A paragraph has line spacing, space before, space after, available bullet formatting, 
#tabs, outline/indentation level, and horizontal alignment. 

#A paragraph can be empty, but if it contains any text, 
#that text is contained in one or more runs.

#A run exists to provide character level formatting, 
#including font typeface, size, and color, an optional hyperlink target URL, 
#bold, italic, and underline styles, strikethrough, kerning, and a few capitalization styles like all caps.

##Accessing the text frame
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame
    # do things with the text frame
    ...

##Accessing paragraphs
#A text frame always contains at least one paragraph. 

#Say for example you want a shape with three paragraphs:

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

text_frame = shape.text_frame
text_frame.clear()  # remove any existing paragraphs, leaving one empty one

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str

##Adding text
#Only runs can actually contain text. 
#Assigning a string to the .text attribute on a shape, text_frame, or paragraph 
#is a shortcut method for placing text in a run contained by those objects. 

shape.text = 'foobar'
# is equivalent to ...
text_frame = shape.text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'foobar'

##Applying text frame-level formatting
#The following produces a shape with a single paragraph, 
#a slightly wider bottom than top margin (these default to 0.05'), no left margin, 
#text aligned top, and word wrapping turned off. 

#In addition, the auto-size behavior is set to adjust the width and height 
#of the shape to fit its text. 
#Note that vertical alignment is set on the text frame. 
#Horizontal alignment is set on each paragraph:

from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

text_frame = shape.text_frame
text_frame.text = 'Spam, eggs, and spam'
text_frame.margin_bottom = Inches(0.08)
text_frame.margin_left = 0
text_frame.vertical_anchor = MSO_ANCHOR.TOP
text_frame.word_wrap = False
text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

##Applying paragraph formatting

#The following produces a shape containing three left-aligned paragraphs, 
#the second and third indented (like sub-bullets) under the first:

from pptx.enum.text import PP_ALIGN

paragraph_strs = [
    'Egg, bacon, sausage and spam.',
    'Spam, bacon, sausage and spam.',
    'Spam, egg, spam, spam, bacon and spam.'
]

text_frame = shape.text_frame
text_frame.clear()

p = text_frame.paragraphs[0]
p.text = paragraph_strs[0]
p.alignment = PP_ALIGN.LEFT

for para_str in paragraph_strs[1:]:
    p = text_frame.add_paragraph()
    p.text = para_str
    p.alignment = PP_ALIGN.LEFT
    p.level = 1

##Applying character formatting
#Character level formatting is applied at the run level, 

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

text_frame = shape.text_frame
text_frame.clear()  # not necessary for newly-created shape

p = text_frame.paragraphs[0]
run = p.add_run()
run.text = 'Spam, eggs, and spam'

font = run.font
font.name = 'Calibri'
font.size = Pt(18)
font.bold = True
font.italic = None  # cause value to be inherited from theme
font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

#can set the font color to an absolute RGB value. 
#Note that this will not change color when the theme is changed:
font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)

#A run can also be made into a hyperlink by providing a target URL:
run.hyperlink.address = 'https://github.com/scanny/python-pptx'

##Working with charts

#The following code adds a single-series column chart in a new presentation:

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

prs.save('chart-01.pptx')

#Customizing things a bit
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))

graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart

#XY and Bubble charts
chart_data = XyChartData()

series_1 = chart_data.add_series('Model 1')
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series('Model 2')
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
).chart

#Creation of a bubble chart 
chart_data = BubbleChartData()

series_1 = chart_data.add_series('Series 1')
series_1.add_data_point(0.7, 2.7, 10)
series_1.add_data_point(1.8, 3.2, 4)
series_1.add_data_point(2.6, 0.8, 8)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
).chart

##Axes in chart 

from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt

category_axis = chart.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)

value_axis = chart.value_axis
value_axis.maximum_scale = 50.0
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"%"'
tick_labels.font.bold = True
tick_labels.font.size = Pt(14)

#Data Labels
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION

plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(13)
data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
data_labels.position = XL_LABEL_POSITION.INSIDE_END

#Legend
from pptx.enum.chart import XL_LEGEND_POSITION

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

#Line Chart
chart_data = ChartData()
chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
chart_data.add_series('West',    (32.2, 28.4, 34.7))
chart_data.add_series('East',    (24.3, 30.6, 20.2))
chart_data.add_series('Midwest', (20.4, 18.3, 26.2))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True

#Pie Chart
#has a single series and doesn't have any axes:

chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


##Working with tables
table
    A table is a matrix of cells arranged in aligned rows and columns. 
cell
    An individual content 'container' within a table. A
    cell has a text-frame in which it holds that content. 
    A PowerPoint table cell can only contain text. 
    It cannot hold images, other shapes, or other tables.
    A cell has a background fill, borders, margins, 
    and several other formatting settings that can be customized on a cell-by-cell basis.
row
    A side-by-side sequence of cells running across the table, 
    all sharing the same top and bottom boundary.
column
    A vertical sequence of cells spanning the height of the table, 
    all sharing the same left and right boundary.
table grid, also cell grid
    The underlying cells in a PowerPoint table are strictly regular. 
    In a three-by-three table there are nine grid cells, 
    three in each row and three in each column. 
    The presence of merged cells can obscure portions of the cell grid, 
    but not change the number of cells in the grid. 
    Access to a table cell in python-pptx is always via that cells coordinates 
    in the cell grid, which may not conform to its visual location 
    (or lack thereof) in the table.
merged cell
    A cell can be 'merged' with adjacent cells, horizontally, vertically, or both, 
    causing the resulting cell to look and behave like a single cell 
    that spans the area formerly occupied by those individual cells.
merge-origin cell
    The top-left grid-cell in a merged cell has certain special behaviors. 
    The content of that cell is what appears on the slide; 
    content of any 'spanned' cells is hidden. 
    In python-pptx a merge-origin cell can be identified 
    with the _Cell.is_merge_origin property. 
    Such a cell can report the size of the merged cell 
    with its span_height and span_width properties, 
    and can be 'unmerged' back to its underlying grid cells 
    using its split() method.
spanned-cell
    A grid-cell other than the merge-origin cell that is 'occupied' 
    by a merged cell is called a spanned cell. 
    Intuitively, the merge-origin cell 'spans' the other grid cells 
    within its area. 
    A spanned cell can be identified with its _Cell.is_spanned property. 
    A merge-origin cell is not itself a spanned cell.

##Adding a table

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
shape = slide.shapes.add_table(3, 3, x, y, cx, cy)

shape   #<pptx.shapes.graphfrm.GraphicFrame object at 0x1022816d0>
shape.has_table #True
table = shape.table #
table   #<pptx.table.Table object at 0x1096f8d90>


##Accessing a cell
cell = table.cell(0, 0)
#a cell has a text-frame and can contain arbitrary text divided into paragraphs and runs
#or with shortcut 
cell.text   #''
cell.text = 'Unladen Swallow'

##Merging cells
#A merged cell is produced by specifying two diagonal cells. 

cell = table.cell(0, 0)
other_cell = table.cell(1, 1)
cell.is_merge_origin #False
cell.merge(other_cell)
cell.is_merge_origin    #True
cell.is_spanned #False
other_cell.is_spanned   #True
table.cell(0, 1).is_spanned #True

##Un-merging a cell
#A merged cell can be restored to its underlying grid cells by calling the split() method on its merge-origin cell. 
#Calling split() on a cell that is not a merge-origin raises ValueError:

cell = table.cell(0, 0)
cell.is_merge_origin    #True
cell.split()
cell.is_merge_origin    #False
table.cell(0, 1).is_spanned #False

##Use Case: Interrogate table for merged cells:

def iter_merge_origins(table):
    """Generate each merge-origin cell in *table*.

    Cell objects are ordered by their position in the table,
    left-to-right, top-to-bottom.
    """
    return (cell for cell in table.iter_cells() if cell.is_merge_origin)

def merged_cell_report(cell):
    """Return str summarizing position and size of merged *cell*."""
    return (
        'merged cell at row %d, col %d, %d cells high and %d cells wide'
        % (cell.row_idx, cell.col_idx, cell.span_height, cell.span_width)
    )

# ---Print a summary line for each merged cell in *table*.---
for merge_origin_cell in iter_merge_origins(table):
    print(merged_cell_report(merge_origin_cell))

#prints a report like:
merged cell at row 0, col 0, 2 cells high and 2 cells wide
merged cell at row 3, col 2, 1 cells high and 2 cells wide
merged cell at row 4, col 0, 2 cells high and 1 cells wide

##Use Case: Access only cells that display text (are not spanned):
def iter_visible_cells(table):
    return (cell for cell in table.iter_cells() if not cell.is_spanned)

##Use Case: Determine whether table contains merged cells:
def has_merged_cells(table):
    for cell in table.iter_cells():
        if cell.is_merge_origin:
            return True
    return False













###.xlsx
$ pip install openpyxl
#https://openpyxl.readthedocs.io/en/stable/api/openpyxl.html

##Write a workbook

from openpyxl import Workbook
from openpyxl.utils import get_column_letter

wb = Workbook()
dest_filename = 'empty_book.xlsx'

ws1 = wb.active
ws1.title = "range names"

#Append complete row, 0-600 for 40 rows 
for row in range(1, 40):
    ws1.append(range(600))

ws2 = wb.create_sheet(title="Pi")
#write directly to F5 (columnRow)
ws2['F5'] = 3.14

ws3 = wb.create_sheet(title="Data")
#Write specifically 
for row in range(10, 20):
    for col in range(27, 54):
        _ = ws3.cell(column=col, row=row, value="{0}".format(get_column_letter(col)))

print(ws3['AA10'].value)

#More writing 
ws = wb.create_sheet(title="Misc")
# set date using a Python datetime
import datetime 
ws['A1'] = datetime.datetime(2010, 7, 21)
ws['A1'].number_format #'yyyy-mm-dd h:mm:ss'
# add a simple formula
ws["A2"] = "=SUM(1, 1)" #or "=SUM(A1:B6)"

#openpyxl never evaluates formula but it is possible to check the name of a formula:
from openpyxl.utils import FORMULAE
"HEX2DEC" in FORMULAE #True 

#Merge / Unmerge cells
#When you merge cells all cells but the top-left one are removed 
ws.merge_cells('A3:D4')
ws.unmerge_cells('A3:D4')
# or equivalently, index starts from 1
ws.merge_cells(start_row=3, start_column=1, end_row=4, end_column=4)
#ws.unmerge_cells(start_row=2, start_column=1, end_row=4, end_column=4)

#Inserting an image
from openpyxl.drawing.image import Image
ws['A5'] = 'You should below'
img = Image('data/monty-truth.png')
# add to worksheet and anchor next to cells
ws.add_image(img, 'A6')

#Line chart 
from datetime import date
from openpyxl.chart import (
    LineChart,
    Reference,
    BarChart,
    PieChart
)
from openpyxl.chart.axis import DateAxis

rows = [
    ['Date', 'Batch 1', 'Batch 2', 'Batch 3'],
    [date(2015,9, 1), 40, 30, 25],
    [date(2015,9, 2), 40, 25, 30],
    [date(2015,9, 3), 50, 30, 45],
    [date(2015,9, 4), 30, 25, 40],
    [date(2015,9, 5), 25, 35, 30],
    [date(2015,9, 6), 20, 40, 35],
]

rows  += [
    ['Pie', 'Sold'],
    ['Apple', 50],
    ['Cherry', 30],
    ['Pumpkin', 10],
    ['Chocolate', 40],
]

ws = wb.create_sheet(title="Charts")
for row in rows:
    ws.append(row) #starts from A1 to D7, other one A8 to B12

#Put some border 
from openpyxl.styles.borders import Border, Side
thin_border = Border(
    left=Side(border_style="thin", color='00000000'),
    right=Side(border_style="thin", color='00000000'),
    top=Side(border_style="thin", color='00000000'),
    bottom=Side(border_style="thin", color='00000000')
)
#from A1 to D7 
#ws.cell(row=3, column=2).border = thin_border
for row in ws['A1': 'B12']:
    for cell in row:
        cell.border = thin_border

for row in ws['C1': 'D7']:
    for cell in row:
        cell.border = thin_border

#Bold header , note style is immutable 
for row in ws['A1': 'D1']:
    for cell in row:
        cell.font = cell.font.copy(bold=True, italic=True)

for row in ws['A8': 'B8']:
    for cell in row:
        cell.font = cell.font.copy(bold=True)


#line chart 
c1 = LineChart()
c1.title = "Sales"
c1.y_axis.title = 'Size'

#crossAx is must (this from excel specification)
#Specifies the axId(axis ID) of axis that this axis cross. 
c1.y_axis.crossAx = 500  #this is for dateAxis
c1.x_axis = DateAxis(crossAx=100) #this is for Numeric Axis , for TextAxis, it is 10, for SeriesAxis, it is 1000
c1.x_axis.number_format = 'd-mmm'
c1.x_axis.majorTimeUnit = "days"
c1.x_axis.title = "Date"
#index starts from 1, data is B1 to D7 including titles
data = Reference(ws, min_col=2, min_row=1, max_col=4, max_row=7)
c1.add_data(data, titles_from_data=True)
#now add date , A2 to A7 
dates = Reference(ws, min_col=1, min_row=2, max_row=7)
c1.set_categories(dates)

#By default the top-left corner of a chart is anchored to cell A15 
#and the size is 15 x 7.5 cm (approximately 5 columns by 14 rows). 
#This can be changed by setting the anchor, width and height properties of the chart.
ws.add_chart(c1, "A15")

#Bar-chart 
#index starts from 1,A8 to B12
labels = Reference(ws, min_col=1, min_row=9, max_row=12)
#labels
#'Charts'!$A$9:$A$12
data = Reference(ws, min_col=2, min_row=8, max_row=12)

c2 = BarChart()
c2.add_data(data, titles_from_data=True)
c2.set_categories(labels)
c2.title = "Pies sold by category"
ws.add_chart(c2, "F15")

#Pie chart 
c3 = PieChart()
c3.add_data(data, titles_from_data=True)
c3.set_categories(labels)
c3.title = "Pies sold by category"
ws.add_chart(c3, "L15")

#now Save 
wb.save(filename = dest_filename)


##Read an existing workbook
#openpyxl does currently not read all possible items in an Excel file 
#so images and charts will be lost from existing files if they are opened and saved with the same name.

from openpyxl import load_workbook

#load_workbook options 
#guess_types will enable or disable (default) type inference when reading cells.
#data_only controls whether cells with formulae have either the formula (default) 
#or the value stored the last time Excel read the sheet.
#keep_vba controls whether any Visual Basic elements are preserved or not (default). 
#If they are preserved they are still not editable.

wb = load_workbook(filename = 'empty_book.xlsx')
sheet_ranges = wb['range names']
print(sheet_ranges['D18'].value) #3

#Sheets properties 
print(wb.get_sheet_names())

active_sheet = wb.active
print(type(active_sheet))

sheet = wb.get_sheet_by_name("Charts")
print(sheet.title)

#Range 
cells = sheet_ranges['A1': 'B6']
for c1, c2 in cells:
    print("{0:8} {1:8}".format(c1.value, c2.value))

#Iterating by rows, ie row by row 
#index from 1
for row in sheet_ranges.iter_rows(min_row=1, min_col=1, max_row=6, max_col=3):
    for cell in row:
        print(cell.value, end=" ")
    print()   

#0 1 2
#0 1 2
#0 1 2
#Iterating by columns ie column by column 
for col in sheet_ranges.iter_cols(min_row=1, min_col=1, max_row=6, max_col=3):
    for cell in col:
        print(cell.value, end=" ")
    print()    

#0 0 0 
#1 1 1 
#2 2 2 

#simple stats 
rows = sheet_ranges.rows

values = []
for row in rows:
    for cell in row:
        values.append(cell.value)

print("Number of values: {0}".format(len(values)))
print("Sum of values: {0}".format(sum(values)))
print("Minimum value: {0}".format(min(values)))
print("Maximum value: {0}".format(max(values)))
print("Mean: {0}".format(stats.mean(values)))
print("Median: {0}".format(stats.median(values)))
print("Standard deviation: {0}".format(stats.stdev(values)))
print("Variance: {0}".format(stats.variance(values)))

#Dimensions
#To get those cells that actually contain data, we can use dimensions.

print(sheet_ranges.dimensions) #A1:WB39
print("Minimum row: {0}".format(sheet_ranges.min_row))
print("Maximum row: {0}".format(sheet_ranges.max_row))
print("Minimum column: {0}".format(sheet_ranges.min_column))
print("Maximum column: {0}".format(sheet_ranges.max_column))

for row in sheet_ranges[sheet_ranges.dimensions]:
    for cell in row:
        print(cell.value)


  


  
##Fold (hiding rows or columns and outline 

import openpyxl
wb = openpyxl.Workbook()
ws = wb.create_sheet()
ws.column_dimensions.group('A','D', hidden=True)
ws.row_dimensions.group(1,10, hidden=True)
wb.save('group.xlsx')

##Freeze A1 
from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.append(range(10))
c = ws['B2']
ws.freeze_panes = c
wb.save('test.xlsx')

##Inserting rows and columns
#openpyxl.worksheet.worksheet.Worksheet.insert_rows(idx, amount=1))
#   Insert row or rows before row==idx
#openpyxl.worksheet.worksheet.Worksheet.insert_cols(idx, amount=1))
#openpyxl.worksheet.worksheet.Worksheet.delete_rows(idx, amount=1))
#openpyxl.worksheet.worksheet.Worksheet.delete_cols(idx, amount=1))

#The default is one row or column. 
#For example to insert a row at 7 (before the existing row 7):
ws.insert_rows(7)

#To delete the columns F:H:
ws.delete_cols(6, 3)

#Moving ranges of cells
#move the cells in the range D4:F10 up one row, and right two columns.
#The cells will overwrite any existing cells.
ws.move_range("D4:F10", rows=-1, cols=2)
#If cells contain formulae you can let openpyxl translate these 
#This will move the relative references in formulae in the range by one row and one column.
ws.move_range("G4:H10", rows=1, cols=1, translate=True)


##Working with Pandas Dataframes - Reading excel 

from openpyxl.utils.dataframe import dataframe_to_rows
import pandas as pd 
import numpy as np 

wb = Workbook()
ws = wb.active

for r in dataframe_to_rows(df, index=True, header=True):
    ws.append(r)

#To convert a dataframe into a worksheet highlighting the header and index:
for cell in ws['A'] + ws[1]:
    cell.style = 'Pandas'

wb.save("pandas_openpyxl.xlsx")



##Converting a worksheet to a Dataframe

#use the values property (no headers or indices)
df = DataFrame(ws.values)

#with headers or indices
import itertools
data = ws.values
cols = next(data)[1:]
data = list(data)
idx = [r[0] for r in data]
data = (itertools.islice(r, 1, None) for r in data)
df = DataFrame(data, index=idx, columns=cols)

##WorkSheet functions 
from openpyxl import Workbook
wb = Workbook()
ws = wb.active

>>> dir(ws)
['BREAK_COLUMN', 'BREAK_NONE', 'BREAK_ROW', 'HeaderFooter', 'ORIENTATION_LANDSCAPE', 
'ORIENTATION_PORTRAIT', 'PAPERSIZE_A3', 'PAPERSIZE_A4', 'PAPERSIZE_A4_SMALL', 
'PAPERSIZE_A5', 'PAPERSIZE_EXECUTIVE', 'PAPERSIZE_LEDGER', 'PAPERSIZE_LEGAL', 
'PAPERSIZE_LETTER', 'PAPERSIZE_LETTER_SMALL', 'PAPERSIZE_STATEMENT', 
'PAPERSIZE_TABLOID', 'SHEETSTATE_HIDDEN', 'SHEETSTATE_VERYHIDDEN', 
'SHEETSTATE_VISIBLE', 'active_cell', 'add_chart', 'add_data_validation', 
'add_image', 'add_pivot', 'add_table', 'append', 'auto_filter', 'calculate_dimension', 
'cell', 'col_breaks','column_dimensions', 'columns', 'conditional_formatting', 
'data_validations', 'delete_cols', 'delete_rows', 'dimensions', 'encoding', 
'evenFooter', 'evenHeader', 'firstFooter', 'firstHeader', 'formula_attributes', 
'freeze_panes', 'insert_cols', 'insert_rows', 'iter_cols', 'iter_rows', 
'legacy_drawing', 'max_column', 'max_row', 'merge_cells', 'merged_cell_ranges', 
'merged_cells', 'mime_type', 'min_column', 'min_row', 'move_range', 'oddFooter', 
'oddHeader', 'orientation', 'page_breaks', 'page_margins', 'page_setup', 
'paper_size', 'parent', 'path', 'print_area', 'print_options', 'print_title_cols', 
'print_title_rows', 'print_titles', 'protection', 'row_breaks', 'row_dimensions', 
'rows', 'scenarios', 'selected_cell', 'set_printer_settings', 'sheet_format', 
'sheet_properties', 'sheet_state', 'sheet_view', 'show_gridlines', 'show_summary_below',
 'show_summary_right', 'title', 'unmerge_cells','values', 'views']

#Accessing 
#The key can be a single cell coordinate 'A1', a range of cells 'A1:D25',
#individual rows or columns eg 'A', 4 
#or ranges of rows or columns 'A:D',4:10.
#Returns either a single cell or a tuple of rows or columns.
for i in range(30):
    ws.append(range(20))

ws['A1'] #<Cell 'Sheet'.A1>
ws['A1'].value #note only cell has .value 

ws['A1': 'D25'] #((<Cell 'Sheet'.A1>, <Cell 'Sheet'.B1>,...),...)
#or 
ws['A1:D25']  #((<Cell 'Sheet'.A1>, <Cell 'Sheet'.B1>,...),...)
[ [c.value  for c in r]   for r in ws['A1': 'D25']]
 
ws['A']
ws[4]  #index starts from 1 #(<Cell 'Sheet'.A4>, <Cell 'Sheet'.B4>...)
ws['A:D']
ws['A':'D']
ws[1:4]

#Others are 
list(ws.values)
list(ws.rows)
list(ws.columns)
 
 
##Creating a chart
#Charts are composed of at least one series of one or more data points. 
#Series themselves are comprised of references to cell ranges.

from openpyxl import Workbook
wb = Workbook()
ws = wb.active
for i in range(10):
    ws.append([i])

from openpyxl.chart import BarChart, Reference, Series
values = Reference(ws, min_col=1, min_row=1, max_col=1, max_row=10)
chart = BarChart()
chart.add_data(values)
#By default the top-left corner of a chart is anchored to cell E15 
#and the size is 15 x 7.5 cm (approximately 5 columns by 14 rows). 
ws.add_chart(chart, "E15")
wb.save("SampleChart.xlsx")


##Bar and Column Charts
#In bar charts values are plotted as either horizontal bars or vertical columns.
#Switch between vertical and horizontal bar charts by setting 'type' to 'col' or 'bar' respectively.
#When using stacked charts the overlap needs to be set to 100.
#If bars are horizontal, x and y axes are revesed.


from openpyxl import Workbook
from openpyxl.chart import BarChart, Series, Reference
from copy import deepcopy

wb = Workbook(write_only=True)
ws = wb.create_sheet()

rows = [
    ('Number', 'Batch 1', 'Batch 2'),
    (2, 10, 30),
    (3, 40, 60),
    (4, 50, 70),
    (5, 20, 10),
    (6, 10, 40),
    (7, 50, 30),
]


for row in rows:
    ws.append(row)


chart1 = BarChart()
chart1.type = "col"  #vertical
chart1.style = 10 #builtin style id 
chart1.title = "Bar Chart"
chart1.y_axis.title = 'Test number'
chart1.x_axis.title = 'Sample length (mm)'

data = Reference(ws, min_col=2, min_row=1, max_row=7, max_col=3)
cats = Reference(ws, min_col=1, min_row=2, max_row=7)
chart1.add_data(data, titles_from_data=True)
chart1.set_categories(cats)
chart1.shape = 4
ws.add_chart(chart1, "A10")


chart2 = deepcopy(chart1)
chart2.style = 11
chart2.type = "bar"
chart2.title = "Horizontal Bar Chart"
ws.add_chart(chart2, "G10")


chart3 = deepcopy(chart1)
chart3.type = "col"
chart3.style = 12
chart3.grouping = "stacked"
chart3.overlap = 100
chart3.title = 'Stacked Chart'
ws.add_chart(chart3, "A27")


chart4 = deepcopy(chart1)
chart4.type = "bar"
chart4.style = 13
chart4.grouping = "percentStacked"
chart4.overlap = 100
chart4.title = 'Percent Stacked Chart'
ws.add_chart(chart4, "G27")

wb.save("bar.xlsx")

##3D Bar Charts
from openpyxl import Workbook
from openpyxl.chart import (
    Reference,
    Series,
    BarChart3D,
)

wb = Workbook()
ws = wb.active

rows = [
    (None, 2013, 2014),
    ("Apples", 5, 4),
    ("Oranges", 6, 2),
    ("Pears", 8, 3)
]

for row in rows:
    ws.append(row)

data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=4)
titles = Reference(ws, min_col=1, min_row=2, max_row=4)
chart = BarChart3D()
chart.title = "3D Bar Chart"
chart.add_data(data=data, titles_from_data=True)
chart.set_categories(titles)

ws.add_chart(chart, "E5")
wb.save("bar3d.xlsx")


##2D Area Charts
#Area charts are similar to line charts with the addition 
#that the area underneath the plotted line is filled. 
#Different variants are available by setting the grouping to 'standard', 'stacked' 
#or 'percentStacked'; 'standard' is the default.

from openpyxl import Workbook
from openpyxl.chart import (
    AreaChart,
    Reference,
    Series,
)

wb = Workbook()
ws = wb.active

rows = [
    ['Number', 'Batch 1', 'Batch 2'],
    [2, 40, 30],
    [3, 40, 25],
    [4, 50, 30],
    [5, 30, 10],
    [6, 25, 5],
    [7, 50, 10],
]

for row in rows:
    ws.append(row)

chart = AreaChart()
chart.title = "Area Chart"
chart.style = 13
chart.x_axis.title = 'Test'
chart.y_axis.title = 'Percentage'

cats = Reference(ws, min_col=1, min_row=1, max_row=7)
data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=7)
chart.add_data(data, titles_from_data=True)
chart.set_categories(cats)

ws.add_chart(chart, "A10")

#3D
from openpyxl import Workbook
from openpyxl.chart import (
    AreaChart3D,
    Reference,
    Series,
)


chart = AreaChart3D()
chart.title = "Area Chart"
chart.style = 13
chart.x_axis.title = 'Test'
chart.y_axis.title = 'Percentage'
chart.legend = None

cats = Reference(ws, min_col=1, min_row=1, max_row=7)
data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=7)
chart.add_data(data, titles_from_data=True)
chart.set_categories(cats)

ws.add_chart(chart, "D10")

wb.save("area.xlsx")

##Scatter Charts
#Scatter, or xy, charts are similar to some line charts. 
#The main difference is that one series of values is plotted against another. 
#This is useful where values are unordered.

from openpyxl import Workbook
from openpyxl.chart import (
    ScatterChart,
    Reference,
    Series,
)

wb = Workbook()
ws = wb.active

rows = [
    ['Size', 'Batch 1', 'Batch 2'],
    [2, 40, 30],
    [3, 40, 25],
    [4, 50, 30],
    [5, 30, 25],
    [6, 25, 35],
    [7, 20, 40],
]

for row in rows:
    ws.append(row)

chart = ScatterChart()
chart.title = "Scatter Chart"
chart.style = 13
chart.x_axis.title = 'Size'
chart.y_axis.title = 'Percentage'

xvalues = Reference(ws, min_col=1, min_row=2, max_row=7)
for i in range(2, 4):
    values = Reference(ws, min_col=i, min_row=1, max_row=7)
    series = Series(values, xvalues, title_from_data=True)
    chart.series.append(series)

ws.add_chart(chart, "A10")

wb.save("scatter.xlsx")


##Bubble Charts
#Bubble charts are similar to scatter charts 
#but use a third dimension to determine the size of the bubbles. 
#Charts can include multiple series.

from openpyxl import Workbook
from openpyxl.chart import Series, Reference, BubbleChart

wb = Workbook()
ws = wb.active

rows = [
    ("Number of Products", "Sales in USD", "Market share"),
    (14, 12200, 15),
    (20, 60000, 33),
    (18, 24400, 10),
    (22, 32000, 42),
    (),
    (12, 8200, 18),
    (15, 50000, 30),
    (19, 22400, 15),
    (25, 25000, 50),
]

for row in rows:
    ws.append(row)

chart = BubbleChart()
chart.style = 18 # use a preset style

# add the first series of data
xvalues = Reference(ws, min_col=1, min_row=2, max_row=5)
yvalues = Reference(ws, min_col=2, min_row=2, max_row=5)
size = Reference(ws, min_col=3, min_row=2, max_row=5)
series = Series(values=yvalues, xvalues=xvalues, zvalues=size, title="2013")
chart.series.append(series)

# add the second
xvalues = Reference(ws, min_col=1, min_row=7, max_row=10)
yvalues = Reference(ws, min_col=2, min_row=7, max_row=10)
size = Reference(ws, min_col=3, min_row=7, max_row=10)
series = Series(values=yvalues, xvalues=xvalues, zvalues=size, title="2014")
chart.series.append(series)

# place the chart starting in cell E1
ws.add_chart(chart, "E1")
wb.save("bubble.xlsx")

##Line Charts
#Line charts allow data to be plotted against a fixed axis. 
# there are three kinds of line charts: standard, stacked and percentStacked.

from datetime import date

from openpyxl import Workbook
from openpyxl.chart import (
    LineChart,
    Reference,
)
from openpyxl.chart.axis import DateAxis

wb = Workbook()
ws = wb.active

rows = [
    ['Date', 'Batch 1', 'Batch 2', 'Batch 3'],
    [date(2015,9, 1), 40, 30, 25],
    [date(2015,9, 2), 40, 25, 30],
    [date(2015,9, 3), 50, 30, 45],
    [date(2015,9, 4), 30, 25, 40],
    [date(2015,9, 5), 25, 35, 30],
    [date(2015,9, 6), 20, 40, 35],
]

for row in rows:
    ws.append(row)

c1 = LineChart()
c1.title = "Line Chart"
c1.style = 13
c1.y_axis.title = 'Size'
c1.x_axis.title = 'Test Number'

data = Reference(ws, min_col=2, min_row=1, max_col=4, max_row=7)
c1.add_data(data, titles_from_data=True)

# Style the lines
s1 = c1.series[0]
s1.marker.symbol = "triangle"
s1.marker.graphicalProperties.solidFill = "FF0000" # Marker filling
s1.marker.graphicalProperties.line.solidFill = "FF0000" # Marker outline

s1.graphicalProperties.line.noFill = True

s2 = c1.series[1]
s2.graphicalProperties.line.solidFill = "00AAAA"
s2.graphicalProperties.line.dashStyle = "sysDot"
s2.graphicalProperties.line.width = 100050 # width in EMUs

s2 = c1.series[2]
s2.smooth = True # Make the line smooth

ws.add_chart(c1, "A10")

from copy import deepcopy
stacked = deepcopy(c1)
stacked.grouping = "stacked"
stacked.title = "Stacked Line Chart"
ws.add_chart(stacked, "A27")

percent_stacked = deepcopy(c1)
percent_stacked.grouping = "percentStacked"
percent_stacked.title = "Percent Stacked Line Chart"
ws.add_chart(percent_stacked, "A44")

# Chart with date axis
c2 = LineChart()
c2.title = "Date Axis"
c2.style = 12
c2.y_axis.title = "Size"
c2.y_axis.crossAx = 500  #these are predefined 
c2.x_axis = DateAxis(crossAx=100)  #these are predefined 
c2.x_axis.number_format = 'd-mmm'
c2.x_axis.majorTimeUnit = "days"
c2.x_axis.title = "Date"

c2.add_data(data, titles_from_data=True)
dates = Reference(ws, min_col=1, min_row=2, max_row=7)
c2.set_categories(dates)

ws.add_chart(c2, "A61")

wb.save("line.xlsx")

##3D Line Charts
#In 3D line charts the third axis is the same as the legend for the series.

from datetime import date

from openpyxl import Workbook
from openpyxl.chart import (
    LineChart3D,
    Reference,
)
from openpyxl.chart.axis import DateAxis

wb = Workbook()
ws = wb.active

rows = [
    ['Date', 'Batch 1', 'Batch 2', 'Batch 3'],
    [date(2015,9, 1), 40, 30, 25],
    [date(2015,9, 2), 40, 25, 30],
    [date(2015,9, 3), 50, 30, 45],
    [date(2015,9, 4), 30, 25, 40],
    [date(2015,9, 5), 25, 35, 30],
    [date(2015,9, 6), 20, 40, 35],
]

for row in rows:
    ws.append(row)

c1 = LineChart3D()
c1.title = "3D Line Chart"
c1.legend = None
c1.style = 15
c1.y_axis.title = 'Size'
c1.x_axis.title = 'Test Number'

data = Reference(ws, min_col=2, min_row=1, max_col=4, max_row=7)
c1.add_data(data, titles_from_data=True)

ws.add_chart(c1, "A10")

wb.save("line3D.xlsx")



##Pie Charts
#Pie charts plot data as slices of a circle 
#with each slice representing the percentage of the whole. 
#Slices are plotted in a clockwise direction with 0° being at the top of the circle. 
#Pie charts can only take a single series of data. 
#The title of the chart will default to being the title of the series.

from openpyxl import Workbook

from openpyxl.chart import (
    PieChart,
    ProjectedPieChart,
    Reference
)
from openpyxl.chart.series import DataPoint

data = [
    ['Pie', 'Sold'],
    ['Apple', 50],
    ['Cherry', 30],
    ['Pumpkin', 10],
    ['Chocolate', 40],
]

wb = Workbook()
ws = wb.active

for row in data:
    ws.append(row)

pie = PieChart()
labels = Reference(ws, min_col=1, min_row=2, max_row=5)
data = Reference(ws, min_col=2, min_row=1, max_row=5)
pie.add_data(data, titles_from_data=True)
pie.set_categories(labels)
pie.title = "Pies sold by category"

# Cut the first slice out of the pie
slice = DataPoint(idx=0, explosion=20)
pie.series[0].data_points = [slice]

ws.add_chart(pie, "D1")

#Projected 
ws = wb.create_sheet(title="Projection")

data = [
    ['Page', 'Views'],
    ['Search', 95],
    ['Products', 4],
    ['Offers', 0.5],
    ['Sales', 0.5],
]

for row in data:
    ws.append(row)

projected_pie = ProjectedPieChart()
projected_pie.type = "pie"
projected_pie.splitType = "val" # split by value
labels = Reference(ws, min_col=1, min_row=2, max_row=5)
data = Reference(ws, min_col=2, min_row=1, max_row=5)
projected_pie.add_data(data, titles_from_data=True)
projected_pie.set_categories(labels)

ws.add_chart(projected_pie, "A10")

from copy import deepcopy
projected_bar = deepcopy(projected_pie)
projected_bar.type = "bar"
projected_bar.splitType = 'pos' # split by position

ws.add_chart(projected_bar, "A27")

wb.save("pie.xlsx")


##3D Pie Charts

from openpyxl import Workbook

from openpyxl.chart import (
    PieChart3D,
    Reference
)

data = [
    ['Pie', 'Sold'],
    ['Apple', 50],
    ['Cherry', 30],
    ['Pumpkin', 10],
    ['Chocolate', 40],
]

wb = Workbook()
ws = wb.active

for row in data:
    ws.append(row)

pie = PieChart3D()
labels = Reference(ws, min_col=1, min_row=2, max_row=5)
data = Reference(ws, min_col=2, min_row=1, max_row=5)
pie.add_data(data, titles_from_data=True)
pie.set_categories(labels)
pie.title = "Pies sold by category"


ws.add_chart(pie, "D1")

wb.save("pie3D.xlsx")


##Minima and Maxima
#Axis minimum and maximum values can be set manually to display specific regions on a chart.

from openpyxl import Workbook
from openpyxl.chart import (
    ScatterChart,
    Reference,
    Series,
)

wb = Workbook()
ws = wb.active

ws.append(['X', '1/X'])
for x in range(-10, 11):
    if x:
        ws.append([x, 1.0 / x])

chart1 = ScatterChart()
chart1.title = "Full Axes"
chart1.x_axis.title = 'x'
chart1.y_axis.title = '1/x'
chart1.legend = None

chart2 = ScatterChart()
chart2.title = "Clipped Axes"
chart2.x_axis.title = 'x'
chart2.y_axis.title = '1/x'
chart2.legend = None

chart2.x_axis.scaling.min = 0
chart2.y_axis.scaling.min = 0
chart2.x_axis.scaling.max = 11
chart2.y_axis.scaling.max = 1.5

x = Reference(ws, min_col=1, min_row=2, max_row=22)
y = Reference(ws, min_col=2, min_row=2, max_row=22)
s = Series(y, xvalues=x)
chart1.append(s)
chart2.append(s)

ws.add_chart(chart1, "C1")
ws.add_chart(chart2, "C15")

wb.save("minmax.xlsx")

##Logarithmic Scaling
#Both the x- and y-axes can be scaled logarithmically. The base of the logarithm can be set to any valid float. If the x-axis is scaled logarithmically, negative values in the domain will be discarded.

from openpyxl import Workbook
from openpyxl.chart import (
    ScatterChart,
    Reference,
    Series,
)
import math

wb = Workbook()
ws = wb.active

ws.append(['X', 'Gaussian'])
for i, x in enumerate(range(-10, 11)):
    ws.append([x, "=EXP(-(($A${row}/6)^2))".format(row = i + 2)])

chart1 = ScatterChart()
chart1.title = "No Scaling"
chart1.x_axis.title = 'x'
chart1.y_axis.title = 'y'
chart1.legend = None

chart2 = ScatterChart()
chart2.title = "X Log Scale"
chart2.x_axis.title = 'x (log10)'
chart2.y_axis.title = 'y'
chart2.legend = None
chart2.x_axis.scaling.logBase = 10

chart3 = ScatterChart()
chart3.title = "Y Log Scale"
chart3.x_axis.title = 'x'
chart3.y_axis.title = 'y (log10)'
chart3.legend = None
chart3.y_axis.scaling.logBase = 10

chart4 = ScatterChart()
chart4.title = "Both Log Scale"
chart4.x_axis.title = 'x (log10)'
chart4.y_axis.title = 'y (log10)'
chart4.legend = None
chart4.x_axis.scaling.logBase = 10
chart4.y_axis.scaling.logBase = 10

chart5 = ScatterChart()
chart5.title = "Log Scale Base e"
chart5.x_axis.title = 'x (ln)'
chart5.y_axis.title = 'y (ln)'
chart5.legend = None
chart5.x_axis.scaling.logBase = math.e
chart5.y_axis.scaling.logBase = math.e

x = Reference(ws, min_col=1, min_row=2, max_row=22)
y = Reference(ws, min_col=2, min_row=2, max_row=22)
s = Series(y, xvalues=x)
chart1.append(s)
chart2.append(s)
chart3.append(s)
chart4.append(s)
chart5.append(s)

ws.add_chart(chart1, "C1")
ws.add_chart(chart2, "I1")
ws.add_chart(chart3, "C15")
ws.add_chart(chart4, "I15")
ws.add_chart(chart5, "F30")

wb.save("log.xlsx")

Axis Orientation

Axes can be displayed “normally” or in reverse. Axis orientation is controlled by the scaling orientation property, which can have a value of either 'minMax' for normal orientation or 'maxMin' for reversed.

from openpyxl import Workbook
from openpyxl.chart import (
    ScatterChart,
    Reference,
    Series,
)

wb = Workbook()
ws = wb.active

ws["A1"] = "Archimedean Spiral"
ws.append(["T", "X", "Y"])
for i, t in enumerate(range(100)):
    ws.append([t / 16.0, "=$A${row}*COS($A${row})".format(row = i + 3),
                         "=$A${row}*SIN($A${row})".format(row = i + 3)])

chart1 = ScatterChart()
chart1.title = "Default Orientation"
chart1.x_axis.title = 'x'
chart1.y_axis.title = 'y'
chart1.legend = None

chart2 = ScatterChart()
chart2.title = "Flip X"
chart2.x_axis.title = 'x'
chart2.y_axis.title = 'y'
chart2.legend = None
chart2.x_axis.scaling.orientation = "maxMin"
chart2.y_axis.scaling.orientation = "minMax"

chart3 = ScatterChart()
chart3.title = "Flip Y"
chart3.x_axis.title = 'x'
chart3.y_axis.title = 'y'
chart3.legend = None
chart3.x_axis.scaling.orientation = "minMax"
chart3.y_axis.scaling.orientation = "maxMin"

chart4 = ScatterChart()
chart4.title = "Flip Both"
chart4.x_axis.title = 'x'
chart4.y_axis.title = 'y'
chart4.legend = None
chart4.x_axis.scaling.orientation = "maxMin"
chart4.y_axis.scaling.orientation = "maxMin"

x = Reference(ws, min_col=2, min_row=2, max_row=102)
y = Reference(ws, min_col=3, min_row=2, max_row=102)
s = Series(y, xvalues=x)
chart1.append(s)
chart2.append(s)
chart3.append(s)
chart4.append(s)

ws.add_chart(chart1, "D1")
ws.add_chart(chart2, "J1")
ws.add_chart(chart3, "D15")
ws.add_chart(chart4, "J15")

wb.save("orientation.xlsx")

##Adding a second axis(ie two charts overlapping)
#Adding a second axis actually involves creating a second chart that shares a common x-axis with the first chart but has a separate y-axis.

from openpyxl import Workbook
from openpyxl.chart import (
    LineChart,
    BarChart,
    Reference,
    Series,
)

wb = Workbook()
ws = wb.active

rows = [
    ['Aliens', 2, 3, 4, 5, 6, 7],
    ['Humans', 10, 40, 50, 20, 10, 50],
]

for row in rows:
    ws.append(row)

c1 = BarChart()
v1 = Reference(ws, min_col=1, min_row=1, max_col=7)
c1.add_data(v1, titles_from_data=True, from_rows=True)

c1.x_axis.title = 'Days'
c1.y_axis.title = 'Aliens'
c1.y_axis.majorGridlines = None
c1.title = 'Survey results'


# Create a second chart
c2 = LineChart()
v2 = Reference(ws, min_col=1, min_row=2, max_col=7)
c2.add_data(v2, titles_from_data=True, from_rows=True)
c2.y_axis.axId = 200
c2.y_axis.title = "Humans"

# Display y-axis of the second chart on the right by setting it to cross the x-axis at its maximum
c1.y_axis.crosses = "max"
c1 += c2

ws.add_chart(c1, "D4")

wb.save("secondary.xlsx")


##Chartsheets
#Chartsheets are special worksheets which only contain charts. 
#All the data for the chart must be on a different worksheet.

from openpyxl import Workbook

from openpyxl.chart import PieChart, Reference, Series

wb = Workbook()
ws = wb.active
cs = wb.create_chartsheet()

rows = [
    ["Bob", 3],
    ["Harry", 2],
    ["James", 4],
]

for row in rows:
    ws.append(row)


chart = PieChart()
labels = Reference(ws, min_col=1, min_row=1, max_row=3)
data = Reference(ws, min_col=2, min_row=1, max_row=3)
chart.series = (Series(data),)
chart.title = "PieChart"

cs.add_chart(chart)

wb.save("demo.xlsx")



##Changing the layout of plot area and legend

#The layout of the chart within the canvas can be set by using the layout property 

#A chart cannot be positioned outside of its container 
#the width and height are the dominant constraints: if x + w > 1, then x = 1 - w.

#The chart can be positioned within its container. 
#x and y adjust position, w and h adjust the size . The units are proportions of the container. 

#x is the horizontal position from the left y is the vertical position from the top 
#h is the height of the chart relative to its container 
#w is the width of the box


#mode for the relevant attribute(xMode, yMode, hMode, wMode) can also be set to either factor or edge. 
#Factor is the default:
#This specifies how to interpret the (x,y,h,w) element for this layout.
#edge	Specifies that the x/y/w/h shall be interpreted as the relative to left/top/right/bottom of the chart element.
#factor	Specifies that the x/y/w/h shall be interpreted as the absolute x/y/width/height of the chart element.
#Note The units are always proportions of the container. 
chart.layout.xMode = edge

#The layoutTarget can be set to outer or inner. The default is outer:
#This element specifies whether to layout the plot area by its inside (not including axis and axis labels) or outside (including axis and axis labels).
chart.layout.layoutTarget = inner

#OR can be with ManualLayout
chart.layout = ManualLayout(layoutTarget=None, 
    xMode=None, yMode=None, wMode='factor', hMode='factor', 
    x=None, y=None, w=None, h=None, extLst=None)
    
#The position of the legend can be controlled either by setting its position: 
#r, l, t, b, and tr, for right, left, top, bottom and top right respectively. 
#The default is r.
chart.legend.position = 'tr'

#or applying a manual layout:
chart.legend.layout = ManualLayout(layoutTarget=None, 
    xMode=None, yMode=None, wMode='factor', hMode='factor', 
    x=None, y=None, w=None, h=None, extLst=None)

#Example 
from openpyxl import Workbook, load_workbook
from openpyxl.chart import ScatterChart, Series, Reference
from openpyxl.chart.layout import Layout, ManualLayout

wb = Workbook()
ws = wb.active

rows = [
    ['Size', 'Batch 1', 'Batch 2'],
    [2, 40, 30],
    [3, 40, 25],
    [4, 50, 30],
    [5, 30, 25],
    [6, 25, 35],
    [7, 20, 40],
]

for row in rows:
    ws.append(row)

ch1 = ScatterChart()
xvalues = Reference(ws, min_col=1, min_row=2, max_row=7)
for i in range(2, 4):
    values = Reference(ws, min_col=i, min_row=1, max_row=7)
    series = Series(values, xvalues, title_from_data=True)
    ch1.series.append(series)


ch1.title = "Default layout"
ch1.style = 13  #style code 
ch1.x_axis.title = 'Size'
ch1.y_axis.title = 'Percentage'
ch1.legend.position = 'r'

ws.add_chart(ch1, "B10")

from copy import deepcopy

# Half-size chart, bottom right
ch2 = deepcopy(ch1)
ch2.title = "Manual chart layout"
ch2.legend.position = "tr"
ch2.layout=Layout(
    manualLayout=ManualLayout(
        x=0.25, y=0.25,
        h=0.5, w=0.5,
    )
)
ws.add_chart(ch2, "H10")

# Half-size chart, centred
ch3 = deepcopy(ch1)
ch3.layout = Layout(
    ManualLayout(
    x=0.25, y=0.25,
    h=0.5, w=0.5,
    xMode="edge",
    yMode="edge",
    )
)
ch3.title = "Manual chart layout, edge mode"
ws.add_chart(ch3, "B27")

# Manually position the legend bottom left
ch4 = deepcopy(ch1)
ch4.title = "Manual legend layout"
ch4.legend.layout = Layout(
    manualLayout=ManualLayout(
        yMode='edge',
        xMode='edge',
        x=0, y=0.9,
        h=0.1, w=0.5
    )
)

ws.add_chart(ch4, "H27")

wb.save("chart_layout.xlsx")



##Adding a comment to a cell
#Openpyxl currently supports the reading and writing of comment text only. 
#Formatting information is lost. 
#Comment dimensions are lost upon reading, but can be written

#Comments have a text attribute and an author attribute, which must both be set

from openpyxl import Workbook
from openpyxl.comments import Comment
wb = Workbook()
ws = wb.active
comment = ws["A1"].comment
comment = Comment('This is the comment text', 'Comment Author')
comment.text #'This is the comment text'
comment.author  #'Comment Author'

#If you assign the same comment to multiple cells then openpyxl will automatically create copies

from openpyxl import Workbook
from openpyxl.comments import Comment
wb=Workbook()
ws=wb.active
comment = Comment("Text", "Author")
ws["A1"].comment = comment
ws["B2"].comment = comment
ws["A1"].comment is comment #True
ws["B2"].comment is comment #False

##Loading and saving comments
#Comments present in a workbook when loaded are stored in the comment attribute of their respective cells automatically. 
#Formatting information of comment such as font size, bold and italics are lost, as are the original dimensions and position of the comment's container box.

#Comments remaining in a workbook when it is saved are automatically saved to the workbook file.
#Comment dimensions can be specified for write-only. 
#Comment dimension are in pixels.

from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.utils import units

wb=Workbook()
ws=wb.active

comment = Comment("Text", "Author")
comment.width = 300
comment.height = 50

ws["A1"].comment = comment

wb.save('commented_book.xlsx')

#If needed, openpyxl.utils.units contains helper functions for converting 
#from other measurements such as mm or points to pixels:

from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.utils import units

wb=Workbook()
ws=wb.active

comment = Comment("Text", "Author")
comment.width = units.points_to_pixels(300)
comment.height = units.points_to_pixels(50)

ws["A1"].comment = comment


##Working with styles
#There are two types of styles: cell styles and named styles, also known as style templates.

#The following are the default values
from openpyxl.styles import PatternFill, Border, Side, Alignment, Protection, Font
font = Font(name='Calibri',
                size=11,
                bold=False,
                italic=False,
                vertAlign=None,
                underline='none',
                strike=False,
                color='FF000000')
fill = PatternFill(fill_type=None,
                start_color='FFFFFFFF',
                end_color='FF000000')
border = Border(left=Side(border_style=None,
                          color='FF000000'),
                right=Side(border_style=None,
                           color='FF000000'),
                top=Side(border_style=None,
                         color='FF000000'),
                bottom=Side(border_style=None,
                            color='FF000000'),
                diagonal=Side(border_style=None,
                              color='FF000000'),
                diagonal_direction=0,
                outline=Side(border_style=None,
                             color='FF000000'),
                vertical=Side(border_style=None,
                              color='FF000000'),
                horizontal=Side(border_style=None,
                               color='FF000000')
               )
alignment=Alignment(horizontal='general',
                    vertical='bottom',
                    text_rotation=0,
                    wrap_text=False,
                    shrink_to_fit=False,
                    indent=0)
number_format = 'General'
protection = Protection(locked=True,
                        hidden=False)


##Cell Styles
#Cell styles are shared between objects 
#and once they have been assigned they cannot be changed.(immutable), Use .copy() method  
from openpyxl.styles import colors
from openpyxl.styles import Font, Color
from openpyxl import Workbook
wb = Workbook()
ws = wb.active

a1 = ws['A1']
d4 = ws['D4']
ft = Font(color=colors.RED)
a1.font = ft
d4.font = ft

#Below is not allowed as it is immutable 
a1.font.italic = True 

# If you want to change the color of a Font, you need to reassign it::
a1.font = Font(color=colors.RED, italic=True) # the change only affects A1
#or 
a1.font = a1.font.copy(color=colors.RED, italic=True) # the change only affects A1


##Copying styles
#Styles can also be copied

from openpyxl.styles import Font
from copy import copy

ft1 = Font(name='Arial', size=14)
ft2 = copy(ft1)
ft2.name = "Tahoma"
ft1.name #'Arial'
ft2.name    #'Tahoma'
ft2.size # copied from the
14.0

##Basic Font Colors
#Colors are usually RGB or aRGB hexvalues. 
#The colors module contains some handy constants

from openpyxl.styles import Font
from openpyxl.styles.colors import RED
font = Font(color=RED)
font = Font(color="FFBB00")

#There is also support for legacy indexed colors as well as themes and tints
from openpyxl.styles.colors import Color
c = Color(indexed=32)
c = Color(theme=6, tint=0.5)

##Applying Styles
#Styles are applied directly to cells

from openpyxl.workbook import Workbook
from openpyxl.styles import Font, Fill
wb = Workbook()
ws = wb.active
c = ws['A1']
c.font = Font(size=12)

#Styles can also applied to columns and rows 
#but note that this applies only to cells created (in Excel) after the file is closed. 

#If you want to apply styles to entire rows and columns then you must apply the style to each cell yourself. 
#This is a restriction of the file format:

col = ws.column_dimensions['A']
col.font = Font(bold=True)
row = ws.row_dimensions[1]
row.font = Font(underline="single")


##Styling Merged Cells

from openpyxl.styles import Border, Side, PatternFill, Font, GradientFill, Alignment
from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.merge_cells('B2:F4')

top_left_cell = ws['B2']
top_left_cell.value = "My Cell"

thin = Side(border_style="thin", color="000000")
double = Side(border_style="double", color="ff0000")

top_left_cell.border = Border(top=double, left=thin, right=thin, bottom=double)
top_left_cell.fill = PatternFill("solid", fgColor="DDDDDD")
top_left_cell.fill = fill = GradientFill(stop=("000000", "FFFFFF"))
top_left_cell.font  = Font(b=True, color="FF0000")
top_left_cell.alignment = Alignment(horizontal="center", vertical="center")

wb.save("styled.xlsx")

##Edit Page Setup

from openpyxl.workbook import Workbook

wb = Workbook()
ws = wb.active

ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
ws.page_setup.paperSize = ws.PAPERSIZE_TABLOID
ws.page_setup.fitToHeight = 0
ws.page_setup.fitToWidth = 1

##Named Styles
#In contrast to Cell Styles, Named Styles are mutable. 
#They make sense when you want to apply formatting to lots of different cells at once. 
#once you have assigned a named style to a cell, 
#additional changes to the style will not affect the cell.

#Once a named style has been registered with a workbook, 
#it can be referred to simply by name.


from openpyxl.styles import NamedStyle, Font, Border, Side
highlight = NamedStyle(name="highlight")
highlight.font = Font(bold=True, size=20)
bd = Side(style='thick', color="000000")
highlight.border = Border(left=bd, top=bd, right=bd, bottom=bd)

#register explicitly
wb.add_named_style(highlight)

#OR But named styles will also be registered automatically the first time they are assigned to a cell:
ws['A1'].style = highlight

#Once registered assign the style using just the name:
ws['D5'].style = 'highlight'


##Using builtin styles
ws['D5'].style = 'Comma'

#the names for these styles are stored in their localised forms. 
#openpyxl will only recognise the English names and only exactly as written here. 
    'Normal' # same as no style
Number formats
    'Comma'
    'Comma [0]'
    'Currency'
    'Currency [0]'
    'Percent'
Informative
    'Calculation'
    'Total'
    'Note'
    'Warning Text'
    'Explanatory Text'
Text styles
    'Title'
    'Headline 1'
    'Headline 2'
    'Headline 3'
    'Headline 4'
    'Hyperlink'
    'Followed Hyperlink'
    'Linked Cell'
Comparisons
    'Input'
    'Output'
    'Check Cell'
    'Good'
    'Bad'
    'Neutral'
Highlights
    'Accent1'
    '20 % - Accent1'
    '40 % - Accent1'
    '60 % - Accent1'
    'Accent2'
    '20 % - Accent2'
    '40 % - Accent2'
    '60 % - Accent2'
    'Accent3'
    '20 % - Accent3'
    '40 % - Accent3'
    '60 % - Accent3'
    'Accent4'
    '20 % - Accent4'
    '40 % - Accent4'
    '60 % - Accent4'
    'Accent5'
    '20 % - Accent5'
    '40 % - Accent5'
    '60 % - Accent5'
    'Accent6'
    '20 % - Accent6'
    '40 % - Accent6'
    '60 % - Accent6'
    'Pandas'


##Additional Worksheet Properties
#most used ones are the 'fitTopage' page setup property 
#and the tabColor that define the background color of the worksheet tab.

Available properties for worksheets
    'enableFormatConditionsCalculation'
    'filterMode'
    'published'
    'syncHorizontal'
    'syncRef'
    'syncVertical'
    'transitionEvaluation'
    'transitionEntry'
    'tabColor'

Available fields for page setup properties
    'autoPageBreaks' 'fitToPage'
Available fields for outlines
    'applyStyles'
    'summaryBelow'
    'summaryRight'
    'showOutlineSymbols'

#see http://msdn.microsoft.com/en-us/library/documentformat.openxml.spreadsheet.sheetproperties%28v=office.14%29.aspx_ for details.

#By default, outline properties are intitialized 
#so you can directly modify each of their 4 attributes, 
#while page setup properties don't. 
#If you want modify the latter, you should first initialize a openpyxl.worksheet.properties.PageSetupProperties object 
#with the required parameters. 
#Once done, they can be directly modified by the routine later if needed.

from openpyxl.workbook import Workbook
from openpyxl.worksheet.properties import WorksheetProperties, PageSetupProperties

wb = Workbook()
ws = wb.active

wsprops = ws.sheet_properties
wsprops.tabColor = "1072BA"
wsprops.filterMode = False
wsprops.pageSetUpPr = PageSetupProperties(fitToPage=True, autoPageBreaks=False)
wsprops.outlinePr.summaryBelow = False
wsprops.outlinePr.applyStyles = True
wsprops.pageSetUpPr.autoPageBreaks = True


##Conditional Formatting
#Excel supports three different types of conditional formatting: 
#builtins, standard and custom. 

#Builtins combine specific rules with predefined styles. 
#Standard conditional formats combine specific rules with custom formatting. 
#custom means  custom formulae for applying custom formats using differential styles.

#The basic syntax for creating a formatting rule is:

from openpyxl.formatting import Rule
from openpyxl.styles import Font, PatternFill, Border
from openpyxl.styles.differential import DifferentialStyle
dxf = DifferentialStyle(font=Font(bold=True), fill=PatternFill(start_color='EE1111', end_color='EE1111'))

#Rule(type, dxfId=None, priority=0, stopIfTrue=None, aboveAverage=None, percent=None, bottom=None, operator=None, text=None, timePeriod=None, rank=None, stdDev=None, equalAverage=None, formula=(), colorScale=None, dataBar=None, iconSet=None, extLst=None, dxf=None)
#type , Value must be one of {'timePeriod', 'uniqueValues', 'containsText', 'dataBar', 'notContainsText', 'aboveAverage', 'notContainsBlanks', 'beginsWith', 'notContainsErrors', 'expression', 'cellIs', 'endsWith', 'containsBlanks', 'top10', 'duplicateValues', 'colorScale', 'containsErrors', 'iconSet'}
rule = Rule(type='cellIs', dxf=dxf, formula=["10"])

#Because the signatures for some rules can be quite verbose 
#there are also some convenience factories for creating them.
#check https://openpyxl.readthedocs.io/en/stable/api/openpyxl.formatting.rule.html


##Builtin formats condidtional formatting
#Builtin formats contain a sequence of formatting settings 
#which combine a type with an integer for comparison. 
#Possible types are: 'num', 'percent', 'max', 'min', 'formula', 'percentile'.

##Builtin formats - ColorScale
#You can have color scales with 2 or 3 colors. 
#2 color scales produce a gradient from one color to another; 
#3 color scales use an additional color for 2 gradients.


from openpyxl.formatting.rule import ColorScale, FormatObject
from openpyxl.styles import Color
first = FormatObject(type='min')
last = FormatObject(type='max')
# colors match the format objects:
colors = [Color('AA0000'), Color('00AA00')]
cs2 = ColorScale(cfvo=[first, last], color=colors)
# a three color scale would extend the sequences
mid = FormatObject(type='num', val=40)
colors.insert(1, Color('00AA00'))
cs3 = ColorScale(cfvo=[first, mid, last], color=colors)
# create a rule with the color scale
from openpyxl.formatting.rule import Rule
rule = Rule(type='colorScale', colorScale=cs3)

#There is a convenience function for creating ColorScale rules

from openpyxl.formatting.rule import ColorScaleRule
rule = ColorScaleRule(start_type='percentile', start_value=10, start_color='FFAA0000',
                      mid_type='percentile', mid_value=50, mid_color='FF0000AA',
                      end_type='percentile', end_value=90, end_color='FF00AA00')

                      
                      
                      
##Builtin formats - IconSet
#Choose from the following set of icons: 
#'3Arrows', '3ArrowsGray', '3Flags', '3TrafficLights1', '3TrafficLights2', '3Signs', '3Symbols', '3Symbols2', '4Arrows', '4ArrowsGray', '4RedToBlack', '4Rating', '4TrafficLights', '5Arrows', '5ArrowsGray', '5Rating', '5Quarters'

from openpyxl.formatting.rule import IconSet, FormatObject
first = FormatObject(type='percent', val=0)
second = FormatObject(type='percent', val=33)
third = FormatObject(type='percent', val=67)
iconset = IconSet(iconSet='3TrafficLights1', cfvo=[first, second, third], showValue=None, percent=None, reverse=None)
# assign the icon set to a rule
from openpyxl.formatting.rule import Rule
rule = Rule(type='iconSet', iconSet=iconset)

#There is a convenience function for creating IconSet rules:

from openpyxl.formatting.rule import IconSetRule
rule = IconSetRule('5Arrows', 'percent', [10, 20, 30, 40, 50], showValue=None, percent=None, reverse=None)

##Builtin formats - DataBar

from openpyxl.formatting.rule import DataBar, FormatObject
first = FormatObject(type='min')
second = FormatObject(type='max')
data_bar = DataBar(cfvo=[first, second], color="638EC6", showValue=None, minLength=None, maxLength=None)
# assign the data bar to a rule
from openpyxl.formatting.rule import Rule
rule = Rule(type='dataBar', dataBar=data_bar)

#There is a convenience function for creating DataBar rules:

from openpyxl.formatting.rule import DataBarRule
rule = DataBarRule(start_type='percentile', start_value=10, end_type='percentile', end_value='90',
                   color="FF638EC6", showValue="None", minLength=None, maxLength=None)

##Standard conditional formats
Average
Percent
Unique or duplicate
Value
Rank

#Code example 
from openpyxl import Workbook
from openpyxl.styles import Color, PatternFill, Font, Border
from openpyxl.styles.differential import DifferentialStyle
from openpyxl.formatting.rule import ColorScaleRule, CellIsRule, FormulaRule

wb = Workbook()
ws = wb.active

# Create fill
redFill = PatternFill(start_color='EE1111',
               end_color='EE1111',
               fill_type='solid')

# Add a two-color scale
# Takes colors in excel 'RRGGBB' style.
ws.conditional_formatting.add('A1:A10',
            ColorScaleRule(start_type='min', start_color='AA0000',
                          end_type='max', end_color='00AA00')
                          )

# Add a three-color scale
ws.conditional_formatting.add('B1:B10',
               ColorScaleRule(start_type='percentile', start_value=10, start_color='AA0000',
                           mid_type='percentile', mid_value=50, mid_color='0000AA',
                           end_type='percentile', end_value=90, end_color='00AA00')
                             )

# Add a conditional formatting based on a cell comparison
# addCellIs(range_string, operator, formula, stopIfTrue, wb, font, border, fill)
# Format if cell is less than 'formula'
#other operator :"greaterThan", "greaterThanOrEqual", "lessThan", "lessThanOrEqual",
#"equal", "!=": "notEqual"
ws.conditional_formatting.add('C2:C10',
            CellIsRule(operator='lessThan', formula=['C$1'], stopIfTrue=True, fill=redFill))

# Format if cell is between 'formula'
#stopIfTrue: In Microsoft Excel, you can apply multiple conditional formatting rules together. 
#But if you want to stop processing the other rules when the first condition is achieved, use stopIfTrue=True
ws.conditional_formatting.add('D2:D10',
            CellIsRule(operator='between', formula=['1','5'], stopIfTrue=True, fill=redFill))

# Format using a formula
ws.conditional_formatting.add('E1:E10',
            FormulaRule(formula=['ISBLANK(E1)'], stopIfTrue=True, fill=redFill))

# Aside from the 2-color and 3-color scales, 
#format rules take fonts, borders and fills for styling:
myFont = Font()
myBorder = Border()
ws.conditional_formatting.add('E1:E10',
            FormulaRule(formula=['E1=0'], font=myFont, border=myBorder, fill=redFill))

# Highlight cells that contain particular text by using a special formula
red_text = Font(color="9C0006")
red_fill = PatternFill(bgColor="FFC7CE")
dxf = DifferentialStyle(font=red_text, fill=red_fill)
rule = Rule(type="containsText", operator="containsText", text="highlight", dxf=dxf)
rule.formula = ['NOT(ISERROR(SEARCH("highlight",A1)))']
ws.conditional_formatting.add('A1:F40', rule)
wb.save("test.xlsx")

#to apply a conditional format to more than one cell, 
#say a row of cells which contain a particular value.

ws.append(['Software', 'Developer', 'Version'])
ws.append(['Excel', 'Microsoft', '2016'])
ws.append(['openpyxl', 'Open source', '2.6'])
ws.append(['OpenOffice', 'Apache', '4.1.4'])
ws.append(['Word', 'Microsoft', '2010'])

#to higlight the rows where the developer is Microsoft. 
#We do this by creating an expression rule 
#and using a formula to identify which rows contain software developed by Microsoft.

red_fill = PatternFill(bgColor="FFC7CE")
dxf = DifferentialStyle(fill=red_fill)
r = Rule(type="expression", dxf=dxf, stopIfTrue=True)
r.formula = ['$A2="Microsoft"']
ws.conditional_formatting.add("A1:C10", r)



##Pivot Tables
#openpyxl provides read-support for pivot tables 
#it should be possible to edit and manipulate existing pivot tables, eg. change their ranges or whether they should update automatically settings.

#As is the case for charts, images and tables 
#there is currently no management API for pivot tables 
#so that client code will have to loop over the _pivots list of a worksheet.


from openpyxl import load_workbook
wb = load_workbook("campaign.xlsx")
ws = wb["Results"]
pivot = ws._pivots[0] # any will do as they share the same cache
pivot.cache.refreshOnload = True

##Print Settings

from openpyxl.workbook import Workbook

wb = Workbook()
ws = wb.active

ws.print_options.horizontalCentered = True
ws.print_options.verticalCentered = True

##Print Settings- Headers and Footers
#Headers and footers use their own formatting language. 
#This is fully supported when writing them but,only partially when reading them. 

#There is support for the font, size and color for a left, centre/center, or right element. 
#Granular control (highlighting individuals words) will require applying control codes manually.

#Also supported are evenHeader and evenFooter as well as firstHeader and firstFooter.
from openpyxl.workbook import Workbook

wb = Workbook()
ws = wb.active

ws.oddHeader.left.text = "Page &[Page] of &N"
ws.oddHeader.left.size = 14
ws.oddHeader.left.font = "Tahoma,Bold"
ws.oddHeader.left.color = "CC3366"

##Print Settings- Add Print Titles
#You can print titles on every page to ensure that the data is properly labelled.

from openpyxl.workbook import Workbook

wb = Workbook()
ws = wb.active

ws.print_title_cols = 'A:B' # the first two cols
ws.print_title_rows = '1:1' # the first row

##Print Settings- Add a Print Area
#You can select a part of a worksheet as the only part that you want to print

from openpyxl.workbook import Workbook

wb = Workbook()
ws = wb.active

ws.print_area = 'A1:F10'


##Using filters and sorts
#Filters and sorts can only be configured by openpyxl 
#but will need to be applied in applications like Excel. 
#This is because they actually rearranges or format cells or rows in the range.

#To add a filter you define a range and then add columns and sort conditions:

from openpyxl import Workbook

wb = Workbook()
ws = wb.active

data = [
    ["Fruit", "Quantity"],
    ["Kiwi", 3],
    ["Grape", 15],
    ["Apple", 30],
    ["Peach", 3],
    ["Pomegranate", 3],
    ["Pear", 3],
    ["Tangerine", 3],
    ["Blueberry", 3],
    ["Mango", 13],
    ["Watermelon", 3],
    ["Blackberry", 3],
    ["Orange", 3],
    ["Raspberry", 3],
    ["Banana", 3]
]

for r in data:
    ws.append(r)

ws.auto_filter.ref = "A1:B15"  #creates filter 
ws.auto_filter.add_filter_column(0, ["Kiwi", "Apple", "Mango"]) #selects only three, but not applied 
ws.auto_filter.add_sort_condition("B2:B15") #sort condition 

wb.save("filtered.xlsx")

##Validating cells
#Data validators can be applied to ranges of cells but are not enforced or evaluated. 
#Ranges do not have to be contiguous: 
#eg. 'A1 B2:B5' is contains A1 and the cells B2 to B5 but not A2 or B2.

#Validations without any cell ranges will be ignored when saving a workbook.

from openpyxl import Workbook
from openpyxl.worksheet.datavalidation import DataValidation

# Create the workbook and worksheet we'll be working with
wb = Workbook()
ws = wb.active

# Create a data-validation object with list validation
dv = DataValidation(type="list", formula1='"Dog,Cat,Bat"', allow_blank=True)

# Optionally set a custom error message
dv.error ='Your entry is not in the list'
dv.errorTitle = 'Invalid Entry'

# Optionally set a custom prompt message
dv.prompt = 'Please select from the list'
dv.promptTitle = 'List Selection'

# Add the data-validation object to the worksheet
ws.add_data_validation(dv)

# Create some cells, and add them to the data-validation object
c1 = ws["A1"]
c1.value = "Dog"
dv.add(c1)
c2 = ws["A2"]
c2.value = "An invalid value"
dv.add(c2)

# Or, apply the validation to a range of cells
dv.add('B1:B1048576') # This is the same as for the whole of column B

# Check with a cell is in the validator
"B4" in dv #True

##Other validation examples

#Any whole number:
dv = DataValidation(type="whole")

#Any whole number above 100:
dv = DataValidation(type="whole",
                    operator="greaterThan",
                    formula1=100)

#Any decimal number:
dv = DataValidation(type="decimal")

#Any decimal number between 0 and 1:
dv = DataValidation(type="decimal",
                    operator="between",
                    formula1=0,
                    formula2=1)

#Any date:
dv = DataValidation(type="date")

#or time:
dv = DataValidation(type="time")

#Any string at most 15 characters:
dv = DataValidation(type="textLength",
                    operator="lessThanOrEqual"),
                    formula1=15)

#Cell range validation:
from openpyxl.utils import quote_sheetname
dv = DataValidation(type="list",
                    formula1="{0}!$B$1:$B$10".format(quote_sheetname(sheetname))
                    )

#Custom rule:
#http://www.contextures.com/xlDataVal07.html
dv = DataValidation(type="custom",
                    formula1"=SOMEFORMULA")


##Defined Names
#Defined names are descriptive text that is used to represents a cell, range of cells, formula, or constant value.
#They might contain a constant, a formula, a single cell reference, 
#a range of cells or multiple ranges of cells across different worksheets. Or all of the above. 

#They are defined globally for a workbook and accessed from there defined_names attribue.


#Accessing a range called 'my_range':
my_range = wb.defined_names['my_range']
# if this contains a range of cells then the destinations attribute is not None
dests = my_range.destinations # returns a generator of (worksheet title, cell range) tuples

cells = []
for title, coord in dests:
    ws = wb[title]
    cells.append(ws[coord])


##Worksheet Tables
#Worksheet tables are references to groups of cells. 
#This makes certain operations such as styling the cells in a table easier.

#By default tables are created with a header from the first row 
#and filters for all the columns.

#Table names must be unique within a workbook 
#and table headers and filter ranges must always contain strings. 
#If this is not the case then Excel may consider the file invalid and remove the table.

from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo

wb = Workbook()
ws = wb.active

data = [
    ['Apples', 10000, 5000, 8000, 6000],
    ['Pears',   2000, 3000, 4000, 5000],
    ['Bananas', 6000, 6000, 6500, 6000],
    ['Oranges',  500,  300,  200,  700],
]

# add column headings. NB. these must be strings
ws.append(["Fruit", "2011", "2012", "2013", "2014"])
for row in data:
    ws.append(row)

#Table(id=1, displayName=None, ref=None, name=None, comment=None, tableType=None, headerRowCount=1, insertRow=None, insertRowShift=None, totalsRowCount=None, totalsRowShown=None, published=None, headerRowDxfId=None, dataDxfId=None, totalsRowDxfId=None, headerRowBorderDxfId=None, tableBorderDxfId=None, totalsRowBorderDxfId=None, headerRowCellStyle=None, dataCellStyle=None, totalsRowCellStyle=None, connectionId=None, autoFilter=None, sortState=None, tableColumns=(), tableStyleInfo=None, extLst=None)
tab = Table(displayName="Table1", ref="A1:E5")

# Add a default style with striped rows and banded columns
style = TableStyleInfo(name="TableStyleMedium9", showFirstColumn=False,
                       showLastColumn=False, showRowStripes=True, showColumnStripes=True)
tab.tableStyleInfo = style
ws.add_table(tab)
wb.save("table.xlsx")


##Parsing Formulas
#openpyxl supports limited parsing of formulas embedded in cells. 
#The openpyxl.formula package contains a Tokenizer class to break formulas into their consitutuent tokens. 


from openpyxl.formula import Tokenizer
tok = Tokenizer("""=IF($A$1,"then True",MAX(DEFAULT_VAL,'Sheet 2'!B1))""")
print("\n".join("%12s%11s%9s" % (t.value, t.type, t.subtype) for t in tok.items))
#Output 
         IF(       FUNC     OPEN
        $A$1    OPERAND    RANGE
           ,        SEP      ARG
 "then True"    OPERAND     TEXT
           ,        SEP      ARG
        MAX(       FUNC     OPEN
 DEFAULT_VAL    OPERAND    RANGE
           ,        SEP      ARG
'Sheet 2'!B1    OPERAND    RANGE
           )       FUNC    CLOSE
           )       FUNC    CLOSE

##Translating formulae from one location to another
#It is possible to translate (in the mathematical sense) formulae 
#from one location to another using the openpyxl.formulas.translate.Translator class. 

#For example, there a range of cells B2:E7 with a sum of each row in column F:

from openpyxl.formula.translate import Translator
ws['F2'] = "=SUM(B2:E2)"
# move the formula one colum to the right
ws['G2'] = Translator("=SUM(B2:E2)", origin="F2").translate_formula("G2")
ws['G2'].value #'=SUM(C2:F2)'


##Workbook Protection
#To prevent other users from viewing hidden worksheets, 
#adding, moving, deleting, or hiding worksheets, and renaming worksheets, 
#you can protect the structure of your workbook with a password. 

#The password can be set using the openpyxl.workbook.protection.WorkbookProtection.workbookPassword() property
wb.security.workbookPassword = '...'
wb.security.lockStructure = True

#Similarly removing change tracking and change history from a shared workbook 
#can be prevented by setting another password. 
#This password can be set using the openpyxl.workbook.protection.WorkbookProtection.revisionsPassword() property
wb.security.revisionsPassword = '...'

#Specific setter functions are provided 
#if you need to set the raw password value without using the default hashing algorithm 

hashed_password = ...
wb.security.set_workbook_password(hashed_password, already_hashed=True)

##Worksheet Protection
#Unlike workbook protection, sheet protection may be enabled with or without using a password. 
#Sheet protection is enabled using the openpxyl.worksheet.protection.SheetProtection.sheet attribute or calling enable() or disable():

ws = wb.active
wb.protection.sheet = True
wb.protection.enable()
wb.protection.disable()

#If no password is specified, 
#users can disable configured sheet protection without specifying a password. 
#Otherwise they must supply a password to change configured protections. 
ws = wb.active
ws.protection.password = '...'



###PDF
#Based on the PostScript language, each PDF file encapsulates 
#a complete description of a fixed-layout flat document, 
#including the text, fonts, vector graphics(drawn on canvas), raster images(jpeg,png etc) and other information needed to display i

#PyPDF2 is a pure-python PDF library capable of splitting, merging together, cropping, and transforming the pages of PDF files. 
#It can also add custom data, viewing options, and passwords to PDF files. 
#It can retrieve text and metadata from PDFs as well as merge entire files together.

$ pip install pypdf2

##Extracting Metadata From PDFs

from PyPDF2 import PdfFileReader

def get_info(path):
    with open(path, 'rb') as f:
        pdf = PdfFileReader(f)
        info = pdf.getDocumentInfo()
        number_of_pages = pdf.getNumPages()
    print(str(info).encode('ascii','ignore').decode())
    author = info.author
    creator = info.creator
    producer = info.producer
    subject = info.subject
    title = info.title
    return (author, creator, producer, subject, title)

path = 'data/example.pdf'
get_info(path)


##Extracting Text From PDFs
#PyPDF2 has limited support for extracting text from PDFs. 
#It doesn't have built-in support for extracting images

#to extract the text from the first page of the PDF 

from PyPDF2 import PdfFileReader

def text_extractor(path):
    with open(path, 'rb') as f:
        pdf = PdfFileReader(f)
        print("No pages:", pdf.getNumPages())
        # get the first page
        page = pdf.getPage(0)
        print(page)
        print('Page type: {}'.format(str(type(page))))
        text = page.extractText()
        return text 

path = 'data/example.pdf'
text_extractor(path)

##Splitting PDFs
#The PyPDF2 package gives the ability to split up a single PDF into multiple ones(each page)


import os.path 
from PyPDF2 import PdfFileReader, PdfFileWriter

def pdf_splitter(path):
    fname = os.path.splitext(os.path.basename(path))[0]
    pdf = PdfFileReader(path)
    for page in range(pdf.getNumPages()):
        pdf_writer = PdfFileWriter()
        pdf_writer.addPage(pdf.getPage(page))
        output_filename = '{}_page_{}.pdf'.format(fname, page+1)
        with open(output_filename, 'wb') as out:
            pdf_writer.write(out)
        print('Created: {}'.format(output_filename))

path = 'data/example.pdf'
pdf_splitter(path)

#Extract Range of  pages
import os.path 
from PyPDF2 import PdfFileReader, PdfFileWriter

def pdf_splitter(path, start=0,length=1):
    fname = os.path.splitext(os.path.basename(path))[0]
    pdf = PdfFileReader(path)
    pdf_writer = PdfFileWriter()
    for page in range(start, start+length):        
        pdf_writer.addPage(pdf.getPage(page))
    output_filename = '{}_page{}-{}.pdf'.format(fname, start+1, start+length)
    with open(output_filename, 'wb') as out:
        pdf_writer.write(out)
    print('Created: {}'.format(output_filename))

path = "closing and opening rank.pdf"
path = 'data/example.pdf'
pdf_splitter(path,0,3)

#Merging Multiple PDFs Together

import glob

from PyPDF2 import PdfFileWriter, PdfFileReader

def merger(output_path, input_paths):
    pdf_writer = PdfFileWriter()
    for path in input_paths:
        pdf_reader = PdfFileReader(path)
        for page in range(pdf_reader.getNumPages()):
            pdf_writer.addPage(pdf_reader.getPage(page))
    with open(output_path, 'wb') as fh:
        pdf_writer.write(fh)

paths = ['data/example.pdf', 'data/example1.pdf']
merger('pdf_merger.pdf', paths)

#OR 



from PyPDF2 import PdfFileMerger

def merger(output_path, input_paths):
    pdf_merger = PdfFileMerger()
    for path in input_paths:
        pdf_merger.append(path)
    with open(output_path, 'wb') as fileobj:
        pdf_merger.write(fileobj)

paths = ['data/example.pdf', 'data/example1.pdf']
merger('pdf_merger.pdf', paths)


#Or use PdfFileMerger.merge
PdfFileMerger.merge(self, position, fileobj, bookmark=None, pages=None, import_bookmarks=True):
        """
        Merges the pages from the given file into the output file at the
        specified page number.
        :param int position: The *page number* to insert this file. File will
            be inserted after the given number.
        :param fileobj: A File Object or an object that supports the standard 
            read and seek methods similar to a File Object. Could also be a
            string representing a path to a PDF file.
        :param str bookmark: Optionally, you may specify a bookmark to be 
            applied at the beginning of the included file by supplying the 
            text of the bookmark.
        :param pages: can be a :ref:`Page Range <page-range>` or a 
        ``(start, stop[, step])`` tuple
            to merge only the specified range of pages from the source
            document into the output document.
        :param bool import_bookmarks: You may prevent the source 
        document's bookmarks from being imported by specifying this as 
       ``False``.
        """

##Rotating Pages
#you must rotate in 90 degrees increments. 
#You can rotate the PDF pages either clockwise or counterclockwise. 


from PyPDF2 import PdfFileWriter, PdfFileReader

def rotator(path):
    pdf_writer = PdfFileWriter()
    pdf_reader = PdfFileReader(path)
    page1 = pdf_reader.getPage(0).rotateClockwise(90)
    pdf_writer.addPage(page1)
    page2 = pdf_reader.getPage(1).rotateCounterClockwise(90)
    pdf_writer.addPage(page2)
    with open('pdf_rotator.pdf', 'wb') as fh:
        pdf_writer.write(fh)

rotator('data/example.pdf')

##Overlaying/Watermarking Pages
#supports merging PDF pages together or overlaying pages on top of each other. 


from PyPDF2 import PdfFileWriter, PdfFileReader

def watermark(input_pdf, output_pdf, watermark_pdf):
    watermark = PdfFileReader(watermark_pdf)
    watermark_page = watermark.getPage(0)
    pdf = PdfFileReader(input_pdf)
    pdf_writer = PdfFileWriter()
    for page in range(pdf.getNumPages()):
        pdf_page = pdf.getPage(page)
        pdf_page.mergePage(watermark_page)
        pdf_writer.addPage(pdf_page)
    with open(output_pdf, 'wb') as fh:
        pdf_writer.write(fh)

watermark(input_pdf='data/example.pdf', 
              output_pdf='watermarked_w9.pdf',
              watermark_pdf='data/watermark.pdf')


##PDF Encryption


from PyPDF2 import PdfFileWriter, PdfFileReader

def encrypt(input_pdf, output_pdf, password):
    pdf_writer = PdfFileWriter()
    pdf_reader = PdfFileReader(input_pdf)
    for page in range(pdf_reader.getNumPages()):
        pdf_writer.addPage(pdf_reader.getPage(page))
    pdf_writer.encrypt(user_pwd=password, owner_pwd=None, use_128bit=True)
    with open(output_pdf, 'wb') as fh:
        pdf_writer.write(fh)

encrypt(input_pdf='data/example.pdf',
            output_pdf='encrypted.pdf',
            password='blowfish')






###ReportLab open-source PDF Toolkit
#Good reference - ReportLab - PDF Processing with Python - Michael Driscoll.
#A graphics canvas API that 'draws' PDF pages
#A charts and widgets library for creating reusable data graphics.
#A page layout engine - PLATYPUS ("Page Layout and TYPography Using Scripts") 
# which builds documents from elements such as headlines, paragraphs, fonts, tables and vector graphics.

$ pip install reportlab

##The ReportLab engineers describe PLATYPUS as having several layers(from highest to lowest level):
DocTemplates 
    the outermost container of your page
PageTemplates 
    specifies the layout of your page
Frames 
    kind of like a sizer in a desktop user interface.
    Basically it provides a region that contains other flowables
Flowables 
    A text or graphic element that can be 'flowed'
    across page boundaries, such as a paragraph of text. This does
    not include footers and headers.
pdfgen.Canvas 
    The lowest level of ReportLab, It will actually receive its instructions
    from one or more of the upper layers and 'paint' your documentaccordingly.

    
##ReportLab - hello_reportlab.py

from reportlab.pdfgen import canvas

c = canvas.Canvas("hello.pdf")
c.drawString(100, 100, "Welcome to Reportlab!")
c.showPage()
c.save()

##Units in ReportLab 
#in ReportLab you can position your elements (text, images, etc) using points.
#lib.units has mm, incg, cm in terms of units 
#also x is from left to right and y is from bottom to up 
inch = 72.0
cm = inch / 2.54
mm = cm * 0.1
pica = 12.0

#code 
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import mm

def coord(x, y, height, unit=1):
    x, y = x * unit, height - y * unit
    return x, y

c = canvas.Canvas("hello.pdf", pagesize=letter)
width, height = letter
c.drawString(*coord(15, 20, height, mm), text="Welcome to Reportlab!")
c.showPage()
c.save()

#or Use Canvas’s bottomUp parameter to zero, such that y is from top 
from reportlab.pdfgen import canvas
from reportlab.lib.units import mm

def coord(x, y, unit=1):
    x, y = x * unit, y * unit
    return x, y

c = canvas.Canvas("hello.pdf", bottomup=0)
c.drawString(*coord(15, 20, mm), text="Welcome to Reportlab!")
c.showPage()
c.save()





##ReportLab - font_demo.py

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def font_demo(my_canvas, fonts):
    pos_y = 750  #start from top 
    for font in fonts:
        my_canvas.setFont(font, 12)
        my_canvas.drawString(30, pos_y, font) #draw each font name 
        pos_y -= 10

if __name__ == '__main__':
    my_canvas = canvas.Canvas("font_demo.pdf",
                              pagesize=letter)
    fonts = my_canvas.getAvailableFonts() #get all fonts 
    font_demo(my_canvas, fonts)
    my_canvas.save()


##Translate & rotate 
#Use the translate(inch, inch) method to set our origin from the bottom left 
#to an inch from the bottom left and an inch up

#canvass rotate method to draw text at different angles
#specify the y coordinate in the negative 
#since the coordinate system is now in a rotated state. 
#If you don't do that, your string will be drawn outside the page’s boundary and you won’t see it.

from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas


def rotate_demo():
    my_canvas = canvas.Canvas("rotated.pdf",
                              pagesize=letter)
    my_canvas.translate(inch, inch)
    my_canvas.setFont('Helvetica', 14)
    my_canvas.drawString(inch, inch, 'Normal')
    my_canvas.line(inch, inch, inch+100, inch)

    my_canvas.rotate(45)
    my_canvas.drawString(inch, -inch, '45 degrees')
    my_canvas.line(inch, inch, inch+100, inch)

    my_canvas.rotate(45)
    my_canvas.drawString(inch, -inch, '90 degrees')
    my_canvas.line(inch, inch, inch+100, inch)

    my_canvas.save()


##String Alignment 
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter


def string_alignment(my_canvas):
    width, height = letter
    
    my_canvas.drawString(80, 700, 'Standard String')
    my_canvas.drawRightString(80, 680, 'Right String')
    
    numbers = [987.15, 42, -1,234.56, (456.78)]
    y = 650
    for number in numbers:
        my_canvas.drawAlignedString(60, y, str(number))
        y -= 20
    
    my_canvas.drawCentredString(width / 2, 550, 'Centered String')
    
    my_canvas.showPage()
    

if __name__ == '__main__':
    my_canvas = canvas.Canvas("string_alignment.pdf")
    string_alignment(my_canvas)
    my_canvas.save()
    
##Drawing 
#Canvas supports many drawing
#calors are 
By specifying red/green/blue values (i.e. values must be betweenzero and one)
By name or
By gray level
#for example 
from reportlab.lib import colors

sample_colors = [colors.aliceblue,
        colors.aquamarine,
        colors.lavender,
        colors.beige,
        colors.chocolate]
        
>>> print(colors.aliceblue)
Color(.941176,.972549,1,1)
#example  
def draw_shapes():
    c = canvas.Canvas("draw_other.pdf")
    c.setStrokeColorRGB(0.2, 0.5, 0.3)
    #rect(self, x, y, width, height, stroke=1, fill=0)
    c.rect(10, 740, 100, 80, stroke=1, fill=0)
    #ellipse(self, x1, y1, x2, y2, stroke=1, fill=0)
    c.ellipse(10, 680, 100, 630, stroke=1, fill=1)
    #wedge(self, x1,y1, x2,y2, startAng, extent, stroke=1, fill=0)
    c.wedge(10, 600, 100, 550, 45, 90, stroke=1, fill=0)
    #circle(self, x_cen, y_cen, r, stroke=1, fill=0)
    c.circle(300, 600, 50)

    my_canvas.setLineWidth(.3)
    start_y = 710
    my_canvas.line(30, start_y, 580, start_y)
    
    my_canvas.setFont('Helvetica', 10)
    x = 30
    grays = [0.0, 0.25, 0.50, 0.75, 1.0]
    for gray in grays:
        my_canvas.setFillGray(gray)
        my_canvas.circle(x, 730, 20, fill=1)
        gray_str = "Gray={gray}".format(gray=gray)
        my_canvas.setFillGray(0.0)
        my_canvas.drawString(x-10, 700, gray_str)
        x += 75
        
    c.save()
        
##Image 
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
def add_image(image_path):
    my_canvas = canvas.Canvas("canvas_image.pdf",
    pagesize=letter)
    my_canvas.drawImage(image_path, 30, 600,
    width=100, height=100)
    my_canvas.save()
if __name__ == '__main__':
    image_path = 'data/python_logo.png'
    add_image(image_path)
    
##Drawing multiline string 
def textobject_demo():
    my_canvas = canvas.Canvas("txt_obj.pdf",
                              pagesize=letter)
    # Create textobject
    textobject = my_canvas.beginText()

    # Set text location (x, y)
    textobject.setTextOrigin(10, 730)

    # Set font face and size
    textobject.setFont('Times-Roman', 12)

    # Write a line of text + carriage return
    textobject.textLine(text='Python rocks!')
    
    #or for multiline 
    #dont trim leading whitepages 
    #textLines(self, stuff, trim=1)
    textobject.textLine(stuff='Python rocks!', trim=0)

    # Change text color
    textobject.setFillColor(colors.red)

    # Write red text
    textobject.textLine(text='Python rocks in red!')

    # Write text to the canvas
    my_canvas.drawText(textobject)

    my_canvas.save()


##ReportLab -  hello_platypus.py
'''When you create an instance of SimpleDocTemplate, it will
contain one or more other PageTemplates that are a description of the
layout of each of the pages of your document

BaseDocTemplate(self, filename,
    pagesize=defaultPageSize,
    pageTemplates=[],
    showBoundary=0,
    leftMargin=inch,
    rightMargin=inch,
    topMargin=inch,
    bottomMargin=inch,
    allowSplitting=1,
    title=None,
    author=None,
    _pageBreakQuick=1,
    encrypt=None)

The filename parameter is the only one required ,pass a string in for this parameter that can be the
PDF’s file name or object that implements a write method, such as StringIO, BytesIO, file or
socket type.

showBoundary parameter that you can use to turn on the boundaries of
the frames in your document

The allowSplitting parameter tells ReportLab that it should try to split
each Flowable across the Frames.

encrypt parameter that defaults to None, that is unencrypted. 
If you pass in a string object, that will be the user’s password to the document. 
Alternatively, you can pass an instance of reportlab.lib.pdfencrypt.StandardEncryption to encrypt
your PDF, which will give you more control over the encryption settings.

'''

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

def hello():
    doc = SimpleDocTemplate("hello_platypus.pdf",
                            pagesize=letter,
                            rightMargin=72,
                            leftMargin=72,
                            topMargin=72,
                            bottomMargin=18)
    styles = getSampleStyleSheet()

    flowables = []

    text = "Hello, I'm a Paragraph"
    para = Paragraph(text, style=styles["Normal"])
    flowables.append(para)

    doc.build(flowables)

if __name__ == '__main__':
    hello()  
        

##ReportLab -  flowable_orientation.py
'''
set the pagesize parameter to equal landscape(letter). 
The landscape function just takes a page size and sets it to landscape

'''

from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
from reportlab.platypus import Frame, PageTemplate, NextPageTemplate, Spacer
'''
The PageBreak class inserts a page break in your document,
the NextPageTemplate tells ReportLab what template to use starting on
the next page while the Spacer class will add space between the flowables
'''

def alternate_orientations():
    doc = SimpleDocTemplate("orientations.pdf",
                            pagesize=letter,
                            rightMargin=72,
                            leftMargin=72,
                            topMargin=72,
                            bottomMargin=18)
    styles = getSampleStyleSheet()
    normal = styles["Normal"]

    margin = 0.5 * inch
    frame = Frame(margin, margin, doc.width, doc.height, id='frame')
    '''
    PageTemplate(id=None,frames=[],onPage=_doNothing, onPageEnd=_doNothing,
                 pagesize=None, autoNextPageTemplate=None,
                 cropBox=None,
                 artBox=None,
                 trimBox=None,
                 bleedBox=None,
                 )
    '''
    portrait_template = PageTemplate(id='portrait',
                                     frames=[frame],
                                     pagesize=letter)
    landscape_template = PageTemplate(id='landscape',
                                      frames=[frame],
                                      pagesize=landscape(letter))
    doc.addPageTemplates([portrait_template, landscape_template])

    story = []
    story.append(Paragraph('This is a page in portrait orientation', normal))

    # Change to landscape orientation
    story.append(NextPageTemplate('landscape'))  #name is from PageTemplate
    story.append(PageBreak())
    story.append(Spacer(inch, 2*inch))
    story.append(Paragraph('This is a page in landscape orientation', normal))

    # Change back to portrait
    story.append(NextPageTemplate('portrait'))
    story.append(PageBreak())
    story.append(Paragraph("Now we're back in portrait mode again", normal))

    doc.build(story)

if __name__ == '__main__':
    alternate_orientations()
    
    
##ReportLab - form_letter.py

import time
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
'''
The Spacer class gives us a convenient way to add space between
paragraphs or other flowables, while the Image class gives us a nice way to
insert images into our document
'''

def form_letter():
    doc = SimpleDocTemplate("form_letter.pdf",
                            pagesize=letter,
                            rightMargin=72,
                            leftMargin=72,
                            topMargin=72,
                            bottomMargin=18)
    '''
    Argument details 
        'pagesize':defaultPageSize,
        'pageTemplates':[],
        'showBoundary':0,
        'leftMargin':inch,
        'rightMargin':inch,
        'topMargin':inch,
        'bottomMargin':inch,
        'allowSplitting':1,
        'title':None,
        'author':None,
        'subject':None,
        'creator':None,
        'producer':None,
        'keywords':[],
        'invariant':None,
        'pageCompression':None,
        '_pageBreakQuick':1,
        'rotation':0,
        '_debug':0,
        'encrypt': None,
        'cropMarks': None,
        'enforceColorSpace': None,
        'displayDocTitle': None,
        'lang': None,
        'initialFontName': None,
        'initialFontSize': None,
        'initialLeading': None,
        'cropBox': None,
        'artBox': None,
        'trimBox': None,
        'bleedBox': None,
        'keepTogetherClass': KeepTogether,
    '''
    flowables = []
    logo = "python_logo.png"
    magName = "Pythonista"
    issueNum = 12
    subPrice = "99.00"
    limitedDate = "03/05/2010"
    freeGift = "tin foil hat"

    formatted_time = time.ctime()
    full_name = "Mike Driscoll"
    address_parts = ["411 State St.", "Waterloo, IA 50158"]

    im = Image(logo, 2*inch, 2*inch)
    flowables.append(im)

    styles = getSampleStyleSheet()
    # Modify the Normal Style
    styles["Normal"].fontSize = 12
    styles["Normal"].leading = 14
    
    # Create a Justify style
    styles.add(ParagraphStyle(name='Justify', alignment=TA_JUSTIFY))    
    '''
    ParagraphStyle defaults 
        'fontName':_baseFontName,
        'fontSize':10,
        'leading':12,
        'leftIndent':0,
        'rightIndent':0,
        'firstLineIndent':0,
        'alignment':TA_LEFT,
        'spaceBefore':0,
        'spaceAfter':0,
        'bulletFontName':_baseFontName,
        'bulletFontSize':10,
        'bulletIndent':0,
        #'bulletColor':black,
        'textColor': black,
        'backColor':None,
        'wordWrap':None,        #None means do nothing special
                                #CJK use Chinese Line breaking
                                #LTR RTL use left to right / right to left
                                #with support from pyfribi2 if available
        'borderWidth': 0,
        'borderPadding': 0,
        'borderColor': None,
        'borderRadius': None,
        'allowWidows': 1,
        'allowOrphans': 0,
        'textTransform':None,   #uppercase lowercase (captitalize not yet) or None or absent
        'endDots':None,         #dots on the last line of left/right justified paras
                                #string or object with text and optional fontName, fontSize, textColor & backColor
                                #dy
        'splitLongWords':1,     #make best efforts to split long words
        'underlineWidth': _baseUnderlineWidth,  #underline width
        'bulletAnchor': 'start',    #where the bullet is anchored ie start, middle, end or numeric
        'justifyLastLine': 0,   #n allow justification on the last line for more than n words 0 means don't bother
        'justifyBreaks': 0,     #justify lines broken with <br/>
        'spaceShrinkage': _spaceShrinkage,  #allow shrinkage of percentage of space to fit on line
        'strikeWidth': _baseStrikeWidth,    #stroke width
        'underlineOffset': _baseUnderlineOffset,    #fraction of fontsize to offset underlines
        'underlineGap': _baseUnderlineGap,      #gap for double/triple underline
        'strikeOffset': _baseStrikeOffset,  #fraction of fontsize to offset strikethrough
        'strikeGap': _baseStrikeGap,        #gap for double/triple strike
        'linkUnderline': _platypus_link_underline,
        #'underlineColor':  None,
        #'strikeColor': None,
        'hyphenationLang': _hyphenationLang,
        #'hyphenationMinWordLength': _hyphenationMinWordLength,
        'embeddedHyphenation': _embeddedHyphenation,
        'uriWasteReduce': _uriWasteReduce,
    '''
    flowables.append(Paragraph(formatted_time, styles["Normal"]))
    flowables.append(Spacer(1, 12))    
    """ 
    Paragraph(text, style, bulletText=None, caseSensitive=1)
        text a string of stuff to go into the paragraph.
        style is a style definition as in reportlab.lib.styles.
        bulletText is an optional bullet defintion.
        caseSensitive set this to 0 if you want the markup tags and their attributes to be case-insensitive.

        This class is a flowable that can format a block of text
        into a paragraph with a given style.

        The paragraph Text can contain XML-like markup including the tags:
        <b> ... </b> - bold
        < u [color="red"] [width="pts"] [offset="pts"]> < /u > - underline
            width and offset can be empty meaning use existing canvas line width
            or with an f/F suffix regarded as a fraction of the font size
        < strike > < /strike > - strike through has the same parameters as underline
        <i> ... </i> - italics
        <u> ... </u> - underline
        <strike> ... </strike> - strike through
        <super> ... </super> - superscript
        <sub> ... </sub> - subscript
        <font name=fontfamily/fontname color=colorname size=float>
        <span name=fontfamily/fontname color=colorname backcolor=colorname size=float style=stylename>
        <onDraw name=callable label="a label"/>
        <index [name="callablecanvasattribute"] label="a label"/>
        <link>link text</link>
            attributes of links
                size/fontSize/uwidth/uoffset=num
                name/face/fontName=name
                fg/textColor/color/ucolor=color
                backcolor/backColor/bgcolor=color
                dest/destination/target/href/link=target
                underline=bool turn on underline
        <a>anchor text</a>
            attributes of anchors
                size/fontSize/uwidth/uoffset=num
                fontName=name
                fg/textColor/color/ucolor=color
                backcolor/backColor/bgcolor=color
                href=href
                underline="yes|no"
        <a name="anchorpoint"/>
        <unichar name="unicode character name"/>
        <unichar value="unicode code point"/>
        <img src="path" width="1in" height="1in" valign="bottom"/>
                width="w%" --> fontSize*w/100   idea from Roberto Alsina
                height="h%" --> linewidth*h/100 <ralsina@netmanagers.com.ar>

        The whole may be surrounded by <para> </para> tags

        The <b> and <i> tags will work for the built-in fonts (Helvetica
        /Times / Courier).  For other fonts you need to register a family
        of 4 fonts using reportlab.pdfbase.pdfmetrics.registerFont; then
        use the addMapping function to tell the library that these 4 fonts
        form a family e.g.
        from reportlab.lib.fonts import addMapping
        addMapping('Vera', 0, 0, 'Vera')    #normal
        addMapping('Vera', 0, 1, 'Vera-Italic')    #italic
        addMapping('Vera', 1, 0, 'Vera-Bold')    #bold
        addMapping('Vera', 1, 1, 'Vera-BoldItalic')    #italic and bold

        It will also be able to handle any MathML specified Greek characters.
    """
    # Create return address
    flowables.append(Paragraph(full_name, styles["Normal"]))
    for part in address_parts:
        flowables.append(Paragraph(part.strip(), styles["Normal"]))
        
    flowables.append(Spacer(1, 12))  #gap of wifth and height 
    ptext = 'Dear {}:'.format(full_name.split()[0].strip())
    flowables.append(Paragraph(ptext, styles["Normal"]))
    flowables.append(Spacer(1, 12))

    ptext = '''
    We would like to welcome you to our subscriber
    base for {magName} Magazine! You will receive {issueNum} issues at
    the excellent introductory price of ${subPrice}. Please respond by
    {limitedDate} to start receiving your subscription and get the
    following free gift: {freeGift}.
    '''.format(magName=magName,
               issueNum=issueNum,
               subPrice=subPrice,
               limitedDate=limitedDate,
               freeGift=freeGift)
    flowables.append(Paragraph(ptext, styles["Justify"]))
    flowables.append(Spacer(1, 12))

    ptext = '''Thank you very much and we look
    forward to serving you.'''

    flowables.append(Paragraph(ptext, styles["Justify"]))
    flowables.append(Spacer(1, 12))
    ptext = 'Sincerely,'
    flowables.append(Paragraph(ptext, styles["Normal"]))
    flowables.append(Spacer(1, 48))
    ptext = 'Ima Sucker'
    flowables.append(Paragraph(ptext, styles["Normal"]))
    flowables.append(Spacer(1, 12))
    doc.build(flowables)

if __name__ == '__main__':
    form_letter()      
        
        
##ReportLab - paragraph_inline_images.py

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet


def paragraph_inline_images():
    doc = SimpleDocTemplate("paragraph_inline_images.pdf",
                            pagesize=letter
                            )
    styles = getSampleStyleSheet()

    flowables = []

    ptext = '''Here is a picture:
    <img src="python_logo.png" width="50" height="50"/> in the
    middle of our text'''
    p = Paragraph(ptext, styles['Normal'])
    flowables.append(p)

    doc.build(flowables)

if __name__ == '__main__':
    paragraph_inline_images()

        
##ReportLab -     paragraph_bullets.py

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet


def paragraph_bullets():
    doc = SimpleDocTemplate("paragraph_bullets.pdf",
                            pagesize=letter
                            )
    styles = getSampleStyleSheet()

    flowables = []

    ptext = "I'm a custom bulletted paragraph"
    para = Paragraph(ptext, style=styles["Normal"], bulletText='-')
    flowables.append(para)

    ptext = "This is a normal bullet"
    para = Paragraph(ptext, style=styles["Normal"], bulletText='•')
    flowables.append(para)

    ptext = "<bullet>&bull;</bullet>This text uses the bullet tag"
    para = Paragraph(ptext, style=styles["Normal"])
    flowables.append(para)

    doc.build(flowables)

if __name__ == '__main__':
    paragraph_bullets()   
        
##ReportLab - paragraph_fonts.py

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

def paragraph_fonts():
    doc = SimpleDocTemplate("paragraph_fonts.pdf",
                            pagesize=letter
                            )
    styles = getSampleStyleSheet()

    flowables = []

    ptext = "<font name=helvetica size=12>Welcome to Reportlab! " \
            "(helvetica)</font>"
    para = Paragraph(ptext, style=styles["Normal"])
    flowables.append(para)

    ptext = "<font face=courier size=14>Welcome to Reportlab! " \
            "(courier)</font>"
    para = Paragraph(ptext, style=styles["Normal"])
    flowables.append(para)

    ptext = "<font name=times-roman size=16>Welcome to Reportlab! " \
            "(times-roman)</font>"
    para = Paragraph(ptext, style=styles["Normal"])
    flowables.append(para)

    doc.build(flowables)

if __name__ == '__main__':
    paragraph_fonts()

        

##ReportLab - absolute_pos_flowable.py
'''
a class named Flowable that is a Python abstract class which has the following methods:
    draw
    wrap
    split (optionally)
Instead the calling draw, user should call drawOn(self, canvas, x, y) which will call draw internally.
The wrapOn(self, canv, aW, aH) calculates how to wrap the text in your paragraph based on that information.
'''

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet


def mixed():
    my_canvas = canvas.Canvas("mixed_flowables.pdf", pagesize=letter)
    styles = getSampleStyleSheet()
    width, height = letter

    text = "Hello, I'm a Paragraph"
    para = Paragraph(text, style=styles["Normal"])
    para.wrapOn(my_canvas, width, height)
    para.drawOn(my_canvas, 20, 760)

    my_canvas.save()

if __name__ == '__main__':
    mixed()
    
##ReportLab - two_column_demo.py

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen.canvas import Canvas
from reportlab.platypus import Paragraph, Frame


def frame_demo():
    my_canvas = Canvas("frame_demo.pdf",
                       pagesize=letter)

    styles = getSampleStyleSheet()
    normal = styles['Normal']
    heading = styles['Heading1']

    flowables = []
    flowables.append(Paragraph('Heading #1', heading))
    flowables.append(Paragraph('Paragraph #1', normal))

    right_flowables = []
    right_flowables.append(Paragraph('Heading #2', heading))
    right_flowables.append(Paragraph('ipsum lorem', normal))
    #A Frame is a container that is itself contained within a PageTemplate
    #Frame(x1, y1, width, height, leftPadding=6, bottomPadding=6,
    #        rightPadding=6, topPadding=6, id=None, showBoundary=0)
    left_frame = Frame(inch, inch, width=3*inch, height=9*inch, showBoundary=1)
    right_frame = Frame(4*inch, inch, width=3*inch, height=9*inch)

    left_frame.addFromList(flowables, my_canvas)
    right_frame.addFromList(right_flowables, my_canvas)

    my_canvas.save()

if __name__ == '__main__':
    frame_demo()            
        
        
        
##ReportLab -   table_with_images.py

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.platypus import Paragraph, Image

def table_with_images():
    doc = SimpleDocTemplate("table_with_images.pdf", pagesize=letter)
    story = []
    styles = getSampleStyleSheet()

    img = Image("python_logo.png", 50, 50)

    ptext = 'This is some <font color=blue size=14>formatted</font> text'
    p = Paragraph(ptext, styles['Normal'])

    data = [['col_{}'.format(x) for x in range(1, 6)],
            [p for x in range(1, 6)],
            [img, img, img, img, img]
            ]

    #TableStyle(cmds=None, parent=None, **kw)
    #The cmds argument is just a list of tuples that define what cell formatting you want to apply.
    #The first element in the tuple is the cell formatting command
    #Thesecond and third elements define the cell coordinates that the formatting
    #will apply to. The coordinates are (column, row),
    #The -1 in the second set of coordinates tells ReportLab that we want the formatting
    #to extend across all the columns from left-to-right. 
    #When you use negative values for cell coordinates, you will basically count backwards from the
    #other end of the table,
    #ie  apply the formatting from (0, 0) to (-1, 0), what you are saying is that you want the formatting to apply to the
    #entire first row of the table
    '''
    Command                 Description
    ALIGNMENT (or ALIGN)    LEFT, RIGHT, CENTRE/CENTER or DECIMAL
    BACKGROUND              The cell's background color
    FONT                    The font name to be applied (optionally can add font size and leading)
    FONTNAME (orFACE)       The font name
    FONTSIZE (or SIZE)      The size of the font in points (leading will likely get out of sync)
    LEADING                 The leading space in points
    TEXTCOLOR               The color name string or (R,G,B) tuple
    LEFTPADDING             The amount of left padding as an integer(default: 6)
    RIGHTPADDING            The amount of right padding as an integer (default: 6)
    BOTTOMPADDING The amount of bottom padding as an integer(default: 3)
    TOPPADDING The amount of top padding as an integer(default: 3)
    COLBACKGROUNDS A list of colors that ReportLab will cycle through
    ROWBACKGROUNDS A list of colors that ReportLab will cycle through
    VALIGN TOP, MIDDLE or the default of BOTTOM

    BACKGROUND command
        It will actually take a ReportLab color from reportlab.lib.colors, 
        a string name or a numeric tuple / list. 
        if tuple, tuple must contain the following information:
        (DIRECTION, startColor, endColor). The DIRECTION element needs to
        be either VERTICAL or HORIZONTAL. This will apply the color as a
        gradient
    The GRID command 
        is actually equivalent to applying the BOX and INNERGRID commands. 
        Basically BOX will apply lines to the outside of the specified cells 
        while INNERGRID will add lines inbetween the cells.
    '''
    tblstyle = TableStyle([('INNERGRID', (0,0), (-1,-1), 0.25, colors.black), #0.25 is line width
                           ('BOX', (0,0), (-1,-1), 0.25, colors.black),
                           ])
    tblstyle2 = TableStyle([('BACKGROUND', (0, 0), (-1, 0),
                            ["HORIZONTAL", colors.red, colors.blue]),
                            ('TEXTCOLOR', (0, 1), (-1, 1), colors.blue)
                            ])                      
    tblstyle3 = TableStyle([('FONT', (0, 1), (-1, 1), 'Helvetica', 24)
                            ])
    tblstyle4 = TableStyle([('FONT', (0, 1), (-1, 1), 'Helvetica'),
                            ('FONTSIZE', (0, 1), (-1, 1), 24)
                            ])
    tblstyle5 = TableStyle([('FONT', (0, 0), (-1, 0), 'Times-Roman'),
                            ('FONT', (0, 1), (-1, 1), 'Helvetica', 24),
                            ('FONT', (0, 2), (-1, 2), 'Courier', 12)
                            ])     
    tblstyle6 = TableStyle(
                        [('LINEABOVE', (0, 0), (-1, 0), 0.5, colors.red),
                        ('LINEBELOW', (0, 0), (-1, 0), 1.5, colors.blue),
                        ('LINEBEFORE', (0, 0), (0, -1), 2.5, colors.orange),
                        ('LINEAFTER', (-1, 0), (-1, -1), 3.5, colors.green),
                        ])            
    tblstyle7 = TableStyle([('INNERGRID', (0,0), (-1,-1), 0.25, colors.black),
                            ('BOX', (0,0), (-1,-1), 0.25, colors.black),
                            ('ALIGN', (0, 0), (0, -1), 'CENTER'), # first    column
                            ('VALIGN', (1, 0), (1, -1), 'MIDDLE'), # second   colu\                    mn
                            ('ALIGN', (2, 0), (2, -1), 'CENTER'), # middle   colu\  mn
                            ('VALIGN', (2, 0), (2, -1), 'MIDDLE'), # middle colu\               mn
                            ('ALIGN', (-1, 0), (-1, -1), 'RIGHT'), # last column
                            ])

    tblstyle8 = TableStyle([('ROWBACKGROUNDS', (0,0), (-1,-1), [colors.gray, colors.white]),
                            ('COLBACKGROUNDS', (0,0), (-1,-1),    [colors.red, colors.white, colors.blue])
                            ])     
    #(SPAN, (begin_col, begin_row), (end_col, end, row))
    #for spanning cells 
    tblstyle9 = TableStyle([('INNERGRID', (0,0), (-1,-1), 0.25, colors.black),
                            ('BOX', (0,0), (-1,-1), 0.25, colors.black),
                            ('SPAN', (0, -1), (1, -1))
                            ])
    #Table(data, colWidths=None, rowHeights=None, style=None, splitByRow=1,
    #repeatRows=0, repeatCols=0, rowSplitRange=None, spaceBefore=None,
    #spaceAfter=None)
    tbl = Table(data)
    tbl2 = Table(data, colWidths=[55 for x in range(5)],
            rowHeights=[45 for x in range(len(data))]
            )
    tbl.setStyle(tblstyle)
    story.append(tbl)

    doc.build(story)

if __name__ == '__main__':
    table_with_images()     
        
##ReportLab -  Complex document  
from reportlab.platypus import *
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.rl_config import defaultPageSize
from reportlab.lib.pagesizes import A4,landscape
from reportlab.lib.units import inch,cm,mm #72.0 ,28.34, 2.83
from reportlab.pdfgen import canvas

PAGE_HEIGHT=defaultPageSize[1] #841.8897637795277
PAGE_WIDTH=defaultPageSize[0]  #595.2755905511812
styles = getSampleStyleSheet()
Title = "Hello world"
pageinfo = "platypus example"

def myFirstPage(canvas, doc):
    canvas.saveState()
    canvas.setFont('Times-Bold',16)
    canvas.drawCentredString(PAGE_WIDTH/2.0, PAGE_HEIGHT-108, Title)
    canvas.setFont('Times-Roman',9)
    canvas.drawString(inch, 0.75 * inch,"First Page / %s" % pageinfo)
    canvas.restoreState()

def myLaterPages(canvas, doc):
    canvas.saveState()
    canvas.setFont('Times-Roman', 9)
    canvas.drawString(inch, 0.75 * inch,"Page %d %s" % (doc.page, pageinfo))
    canvas.restoreState()

def go():
    doc = SimpleDocTemplate("myreport.pdf",
                        leftMargin = 0.75*inch,
                        rightMargin = 0.75*inch,
                        topMargin = 1*inch,
                        bottomMargin = 1*inch)
    Story = [Spacer(1,2*inch)]  #width, height,
    style = styles["Normal"]
    for i in range(10):
        bogustext = ("Paragraph number <b>%s</b> " % i) *20
        p = Paragraph(bogustext, style)
        Story.append(p)
        Story.append(Spacer(1,0.2*inch))
    Story.append(PageBreak())
    NormalStyle = tables.TableStyle([
        ('BOX',(0,0),(-1,-1),0.45,colors.blue),
        ('ALIGN',(0,0),(-1,-1),'LEFT'),
        ('BACKGROUND',(0,0),(-1,-1),colors.lightblue)
        ])
    mytable = tables.Table([('test','test'),('test2','test2'),('test3','test3')],
                colWidths = 1* inch,rowHeights= 0.25 *inch,style = NormalStyle)
    Story.append(mytable)    
    doc.build(Story, onFirstPage=myFirstPage, onLaterPages=myLaterPages)

if __name__ == "__main__":
    go()      
        
##ReportLab -  pie_chart_with_legend.py

from reportlab.lib.validators import Auto
from reportlab.graphics.charts.legends import Legend
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.shapes import Drawing, String


def pie_chart_with_legend():
    data = list(range(15, 105, 15))
    drawing = Drawing(width=400, height=200)
    my_title = String(170, 40, 'My Pie Chart', fontSize=14)
    pie = Pie()
    pie.sideLabels = True

    pie.x = 150
    pie.y = 65

    pie.data = data
    pie.labels = [letter for letter in 'abcdefg']
    pie.slices.strokeWidth = 0.5
    drawing.add(my_title)
    drawing.add(pie)
    add_legend(drawing, pie, data)
    drawing.save(formats=['pdf'], outDir='.',
                 fnRoot='pie_chart_with_legend')


def add_legend(draw_obj, chart, data):
    legend = Legend()
    legend.alignment = 'right'
    legend.x = 10
    legend.y = 70
    legend.colorNamePairs = Auto(obj=chart)
    draw_obj.add(legend)


if __name__ == '__main__':
    pie_chart_with_legend()


        
##ReportLab -       line_plot_demo.py

from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.lineplots import LinePlot
from reportlab.graphics.widgets.markers import makeMarker


def line_plot_demo():
    d = Drawing(400, 400)
    line = LinePlot()
    line.x = 50
    line.y = 85
    line.height = 150
    line.width = 250
    line.lineLabelFormat = '%2.0f'

    data = [
            ((1,1), (2,2), (2.5,1), (3,3), (4,5)),
            ((1,2), (2,3), (2.5,2), (3.5,5), (4,6))
        ]
    line.data = data

    line.lines[0].strokeColor = colors.green
    line.lines[1].strokeColor = colors.blue
    line.lines[0].strokeWidth = 3

    line.lines[0].symbol = makeMarker('Circle')
    line.lines[1].symbol = makeMarker('Hexagon')

    line.xValueAxis.valueMin = 0
    line.xValueAxis.valueMax = 10
    line.xValueAxis.valueSteps = [1, 2, 4]
    line.xValueAxis.labelTextFormat = '%2.1f'

    line.yValueAxis.valueMin = 0
    line.yValueAxis.valueMax = 12

    d.add(line, '')
    d.save(formats=['pdf'], outDir='.', fnRoot='line_plot_demo')


if __name__ == '__main__':
    line_plot_demo() 
        
##ReportLab - simple_line_chart.py

from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.platypus import SimpleDocTemplate


def simple_line_chart():
    d = Drawing(280, 250)
    line = HorizontalLineChart()
    line.x = 50
    line.y = 85
    line.height = 150
    line.width = 250
    
    data = [[1, 2, 3, None, None, None, 5],
            [10, 5, 2, 6, 8, 3, 5]
            ]
    line.data = data
    line.categoryAxis.categoryNames = [
        'Dogs', 'Cats', 'Mice', 'Hamsters',
        'Parakeets', 'Gerbils', 'Fish'
    ]

    line.lines[0].strokeColor = colors.green
    line.lines[1].strokeColor = colors.blue
    line.lines[0].strokeWidth = 3
    line.categoryAxis.labels.angle = 45
    line.categoryAxis.labels.dy = -15
    
    d.add(line, '')

    doc = SimpleDocTemplate('simple_line_chart.pdf')
    story = []
    story.append(d)
    doc.build(story)

if __name__ == '__main__':
    simple_line_chart()


        
##ReportLab -     simple_bar_chart.py

from reportlab.lib.colors import PCMYKColor
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart


def simple_bar_chart():
    d = Drawing(280, 250)
    bar = VerticalBarChart()
    bar.x = 50
    bar.y = 85
    data = [[1, 2, 3, None, None, None, 5],
            [10, 5, 2, 6, 8, 3, 5],
            [5, 7, 2, 8, 8, 2, 5],
            [2, 10, 2, 1, 8, 9, 5],
            ]
    bar.data = data
    bar.categoryAxis.categoryNames = ['Year1', 'Year2', 'Year3',
                                      'Year4', 'Year5', 'Year6',
                                      'Year7']

    bar.bars[0].fillColor = PCMYKColor(0,100,100,40,alpha=85)
    bar.bars[1].fillColor = PCMYKColor(23,51,0,4,alpha=85)
    
    d.add(bar, '')

    d.save(formats=['pdf'], outDir='.', fnRoot='simple_bar_chart')

if __name__ == '__main__':
    simple_bar_chart()   
        

        
##ReportLab -    simple_stacked_bar_chart.py

from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.platypus import SimpleDocTemplate


def simple_stacked_bar_chart():
    """
    Creates a bar chart in a PDF
    """
    d = Drawing(280, 250)
    bar = VerticalBarChart()
    bar.x = 50
    bar.y = 85
    data = [[1, 2, 3, None, None, None, 5],
            [10, 5, 2, 6, 8, 3, 5]
            ]
    bar.data = data
    bar.categoryAxis.categoryNames = ['Year1', 'Year2', 'Year3',
                                      'Year4', 'Year5', 'Year6',
                                      'Year7']

    bar.bars[0].fillColor = colors.green
    bar.bars[1].fillColor = colors.blue
    bar.categoryAxis.labels.angle = 45
    bar.categoryAxis.labels.dy = -15
    bar.categoryAxis.style = 'stacked'
    
    d.add(bar, '')

    doc = SimpleDocTemplate('simple_stacked_bar_chart.pdf')
    story = []
    story.append(d)
    doc.build(story)

if __name__ == '__main__':
    simple_stacked_bar_chart()    
        
##ReportLab - simple_horizontal_bar_chart.py

from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import HorizontalBarChart
from reportlab.platypus import SimpleDocTemplate


def simple_horizontal_bar_chart():
    d = Drawing(280, 250)
    bar = HorizontalBarChart()
    bar.x = 50
    bar.y = 85
    bar.height = 225
    bar.width = 250
    data = [[1, 2, 3, None, None],
            [10, 5, 2, 6, 8],
            [5, 7, 2, 8, 8],
            [2, 10, 2, 1, 8],
            ]
    bar.data = data
    bar.categoryAxis.categoryNames = ['Year1', 'Year2', 'Year3',
                                      'Year4', 'Year5', 'Year6',
                                      'Year7']

    bar.bars[0].fillColor = colors.green
    bar.bars[1].fillColor = colors.blue
    bar.bars[2].fillColor = colors.red
    bar.bars[3].fillColor = colors.purple
    
    bar.categoryAxis.labels.angle = 45
    bar.categoryAxis.labels.dx = -15
    
    d.add(bar, '')

    doc = SimpleDocTemplate('simple_horizontal_bar_chart.pdf')
    story = []
    story.append(d)
    doc.build(story)

if __name__ == '__main__':
    simple_horizontal_bar_chart()

##ReportLab - 3D vertical bar 
from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.barcharts import VerticalBarChart3D
from reportlab.platypus import SimpleDocTemplate, PageBreak


def simple_vertical_3d_bar_chart():
    d = Drawing(280, 250)
    bar = VerticalBarChart3D()
    bar.x = 50
    bar.y = 85
    bar.height = 225
    bar.width = 350
    data = [[1, 2, 3, None, None],
            [10, 5, 2, 6, 8],
            [5, 7, 2, 8, 8],
            [2, 10, 2, 1, 8],
            ]
    bar.data = data
    bar.categoryAxis.categoryNames = ['Year1', 'Year2', 'Year3',
                                      'Year4', 'Year5', 'Year6',
                                      'Year7']
    bar.bars[0].fillColor = colors.green
    bar.bars[1].fillColor = colors.blue
    bar.bars[2].fillColor = colors.red
    bar.bars[3].fillColor = colors.purple    
    bar.categoryAxis.labels.angle = 45
    bar.categoryAxis.labels.dy = -15    
    d.add(bar, '')
    doc = SimpleDocTemplate('simple_vertical_3d_bar_chart.pdf')
    story = []
    story.append(d)
    story.append(PageBreak())
    story.append(d)
    doc.build(story)

if __name__ == '__main__':
    simple_vertical_3d_bar_chart()


    
    

    
    
###PyMuPDF
#MuPDF can access files in PDF, XPS, OpenXPS, CBZ, EPUB and FB2 (e-books) formats, 
#and it is known for its top performance and high rendering quality.

$ pip install PyMuPDF


##How to Convert Images
#Input Formats 	Output Formats 	Description
JPEG 	        - 	            Joint Photographic Experts Group
BMP 	        - 	            Windows Bitmap
JXR 	        - 	            JPEG Extended Range
JPX 	        - 	            JPEG 2000
GIF 	        - 	            Graphics Interchange Format
TIFF 	        - 	            Tagged Image File Format
PNG 	        PNG 	        Portable Network Graphics
PNM 	        PNM 	        Portable Anymap
PGM 	        PGM 	        Portable Graymap
PBM 	        PBM 	        Portable Bitmap
PPM 	        PPM 	        Portable Pixmap
PAM 	        PAM 	        Portable Arbitrary Map
TGA 	        TGA 	        Targa Image File
- 	            PSD 	        Adobe Photoshop Document
- 	            PS 	            Adobe Postscript

#Example 
import fitz
pix = fitz.Pixmap("input.xxx")      # input.xxx: a file in any of the supported input formats
pix.writeImage("output.yyy")        # yyy is any of the supported output formats

#The argument of fitz.Pixmap(arg) can be a file or a bytes object containing a file image
#Instead of creating an output file like above, you can also create a bytes object via pix.getImageData("yyy") and pass this around.
 
#Example: Convert JPEG to Photoshop

import fitz

pix = fitz.Pixmap("myfamily.jpg")
pix.writeImage("myfamily.psd")

#Example: Output JPEG via PIL/Pillow

from PIL import Image
import fitz

pix = fitz.Pixmap(...)
img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
img.save("output.jpg", "JPEG")

#Example: Convert JPEG to Tkinter PhotoImage

import fitz
if str is bytes:          # this is Python 2!
    import Tkinter as tk
else:                     # Python 3 or later!
    import tkinter as tk

pix = fitz.Pixmap("input.jpg")
tkimg = tk.PhotoImage(data=pix.getImageData("ppm")) # PPM is among the tk-supported formats


##How to Extract all Document Text
import sys, fitz                            # import the bindings
fname = sys.argv[1]                         # get document filename
doc = fitz.open(fname)                      # open document
out = open(fname + ".txt", "wb")            # open text output
for page in doc:                            # iterate the document pages
    text = page.getText().encode("utf8")    # get plain text (is in UTF-8)
    out.write(text)                         # write text of page
    out.write(b"\n-----\n")                 # write page delimiter
out.close()

#details 
ext = page.getText("type")
#Use one of the following strings for "type" 
1."text": (default) plain text with line breaks. No formatting, no text position details, no images.
2."html": creates a full visual version of the page including any images. This can be displayed with your internet browser.
3."dict": same information level as HTML, but provided as a Python dictionary. See TextPage.extractDICT() for details of its structure.
4."rawdict": a super-set of TextPage.extractDICT(). It additionally provides character detail information like XML. See TextPage.extractRAWDICT() for details of its structure.
5."xhtml": text information level as the TEXT version but includes images. Can also be displayed by internet browsers.
6."xml": contains no images, but full position and font information down to each single text character. Use an XML module to interpret.

##Searching for Text
#A text page consists of blocks (= roughly paragraphs).
#A block consists of either lines and their characters, or an image.
#A line consists of spans.
#A span consists of font information and characters that share a common baseline.

<page>
    <text block>
        <line>
            <span>
                <char>
    <image block>
        <img>


#You can find out, exactly where on a page a certain text string appears:
areas = page.searchFor("mupdf", hit_max = 16)

#This returns  a list of up to 16 rectangles (Rect), 
#each of which surrounds one occurrence of the string 'mupdf' (case insensitive). 

def search_for(text, words):
    """ Search for text in items of list of words
    Args:
        text: string to be searched for
        words: list of items in format delivered by 'getTextWords()'.
            Each word is represented by:
            [x0, y0, x1, y1, word, bno, lno, wno]
            The first 4 entries are the word's rectangle coordinates, the last 3 are just
            technical info (block number, line number, word number).
    Returns:
        List of rectangles, one for each found locations.
    """
    rect_list = []
    for w in words:
        if text in w[4]:
            rect_list.append(fitz.Rect(w[:4]))
    return rect_list

##How to extract text from a rectangle

from operator import itemgetter 
from itertools import groupby
import fitz
doc = fitz.open("<some.file>")     # any supported document type
page = doc[pno]                    # we want text from this page

"""
-------------------------------------------------------------------------------
Identify the rectangle. We use the text search function here. The two
search strings are chosen to be unique, to make our case work.
The two returned rectangle lists both have only one item.
-------------------------------------------------------------------------------
"""
rl1 = page.searchFor("Die Altersübereinstimmung") # rect list one
rl2 = page.searchFor("Bombardement durch.")       # rect list two
rect = rl1[0] | rl2[0]       # union rectangle
# Now we have the rectangle ---------------------------------------------------

"""
Get all words on page in a list of lists. Each word is represented by:
[x0, y0, x1, y1, word, bno, lno, wno]
The first 4 entries are the word's rectangle coordinates, the last 3 are just
technical info (block number, line number, word number).
"""
words = page.getTextWords()
# We subselect from above list.

# Case 1: select the words fully contained in rect
#------------------------------------------------------------------------------
mywords = [w for w in words if fitz.Rect(w[:4]) in rect]
mywords.sort(key = itemgetter(3, 0))   # sort by y1, x0 of the word rect
group = groupby(mywords, key = itemgetter(3))
print("Select the words strictly contained in rectangle")
print("------------------------------------------------")
for y1, gwords in group:
    print(" ".join(w[4] for w in gwords))

# Case 2: select the words which at least intersect the rect
#------------------------------------------------------------------------------
mywords = [w for w in words if fitz.Rect(w[:4]).intersects(rect)]
mywords.sort(key = itemgetter(3, 0))
group = groupby(mywords, key = itemgetter(3))
print("\nSelect the words intersecting the rectangle")
print("-------------------------------------------")
for y1, gwords in group:
    print(" ".join(w[4] for w in gwords))
    
    
##How to Iterate through the xref Table
#A PDF’s xref table is a list of all objects defined in the file. 
#This table may easily contain many thousand entries – 
#https://pymupdf.readthedocs.io/en/latest/app4/#adobemanual for example has over 330000 objects. 
#Table entry '0' is reserved and must not be touched. 
#The following script loops through the xref table and prints each object's definition:

>>> xreflen = doc._getXrefLength()  # number of objects in file
>>> for xref in range(1, xreflen):  # skip item 0!
        print("")
        print("object %i (stream: %s)" % (xref, doc.isStream(xref)))
        print(doc._getXrefString(i, compressed=False))

This produces the following output:

object 1 (stream: False)
<<
    /ModDate (D:20170314122233-04'00')
    /PXCViewerInfo (PDF-XChange Viewer;2.5.312.1;Feb  9 2015;12:00:06;D:20170314122233-04'00')
>>

object 2 (stream: False)
<<
    /Type /Catalog
    /Pages 3 0 R
>>

object 3 (stream: False)
<<
    /Kids [ 4 0 R 5 0 R ]
    /Type /Pages
    /Count 2
>>

object 4 (stream: False)
<<
    /Type /Page
    /Annots [ 6 0 R ]
    /Parent 3 0 R
    /Contents 7 0 R
    /MediaBox [ 0 0 595 842 ]
    /Resources 8 0 R
>>
...
object 7 (stream: True)
<<
    /Length 494
    /Filter /FlateDecode
>>
...

##How to Handle Page Contents
#A PDF page can have one or more contents objects 
#there are also situations when a single contents object is beneficial
#: it is easier to interpret and better compressible than multiple smaller ones.
#Here are two ways of combining multiple contents of a page:

# method 1: use the clean function
for i in range(len(doc)):
    doc[i]._cleanContents() # cleans and combines multiple Contents
    page = doc[i]           # re-read the page (has only 1 contents now)
    cont = page._getContents()[0]
    # do something with the cleaned, combined contents

# method 2: self-concatenate multiple contents
for page in doc:
    cont = b""              # initialize contents
    for xref in page._getContents(): # loop through content xrefs
        cont += doc._getXrefStream(xref)
    # do something with the combined contents




    

    
    
    
##How to Make Images from Document Pages
import sys, fitz                            # import the binding
fname = sys.argv[1]                         # get filename from command line
doc = fitz.open(fname)                      # open document
for page in doc:                            # iterate through the pages
    pix = page.getPixmap(alpha = False)     # render page to an image
    pix.writePNG("page-%i.png" % page.number)    # store image as a PNG

    
##How to Extract Images: PDF Documents

#embedded images are identified by a cross reference number (xref, an integer). 
#If you know this number, you have two ways to access the image's data. 
#The following assumes you have opened a PDF under the name 'doc':
#    Create a Pixmap of the image with instruction pix = fitz.Pixmap(doc, xref). 
     This method is very fast (single digit micro-seconds). 
     The pixmap’s properties (width, height, …) will reflect the ones of the image. 
     As usual, you can save it as a PNG via method Pixmap.writePNG() 
     (or get the corresponding binary data Pixmap.getPNGData()). 
     There is no way to tell which image format the embedded original has.
#    Extract the image with instruction img = doc.extractImage(xref). 
     This is a dictionary containing the binary image data as img["image"]. 
     A number of meta data are also provided – 
     mostly the same as you would find in the pixmap of the image. 
     The major difference is string img["ext"], which specifies the image format: 
     apart from 'png', strings like 'jpeg', 'bmp', 'tiff', etc. can also occur. 
     Use this string as the file extension if you want to store the image. 
     The execution speed of this method should be compared to the combined speed of the statements 
     pix = fitz.Pixmap(doc, xref);pix.getPNGData(). 
     If the embedded image is in PNG format, 
     the speed of Document.extractImage() is about the same 
     (and the binary image data are identical). 
     Otherwise, this method is thousands of times faster, 
     and the image data is much smaller.

#How do I know those cross reference numbers 'xref' of images?'. 
#There are two answers to this:
'Inspect the page objects' 
    Loop through the document's page number list and execute Document.getPageImageList()
    for each page number. 
    The result is a list of list, and its items look like [xref, smask, ...], 
    containing the xref of an image shown on that page. 
    This xref can then be used with one of the above methods. 
    Use this method for valid (undamaged) documents. 
    Be wary however, that the same image may be referenced multiple times 
    (by different pages), so you might want to provide a mechanism avoiding 
    multiple extracts.
    Script: PyMuPDF-Utilities-master.zip    extract-imga.py
'No need to know' 
    Loop through the list of all xrefs of the document 
    and perform a Document.extractImage() for each one. 
    If the returned dictionary is empty, then continue – this xref is no image. 
    Use this method if the PDF is damaged (unusable pages). 
    Note that a PDF often contains 'pseudo-images' ('stencil masks') 
    with the special purpose to specify the transparency of some other image. 
    You may want to provide logic to exclude those from extraction.
    Script: PyMuPDF-Utilities-master.zip  extract-imgb.py


##How to Make one PDF of all your Pictures
#Method 1: Inserting Images as Pages

#The first one converts each image to a PDF page with the same dimensions:

import os, fitz
import PySimpleGUI as psg                    # for showing progress bar
doc = fitz.open()                            # PDF with the pictures
imgdir = "D:/2012_10_05"                     # where the pics are
imglist = os.listdir(imgdir)                 # list of them
imgcount = len(imglist)                      # pic count

for i, f in enumerate(imglist):
    img = fitz.open(os.path.join(imgdir, f)) # open pic as document
    rect = img[0].rect                       # pic dimension
    pdfbytes = img.convertToPDF()            # make a PDF stream
    img.close()                              # no longer needed
    imgPDF = fitz.open("pdf", pdfbytes)      # open stream as PDF
    page = doc.newPage(width = rect.width,   # new page with ...
                       height = rect.height) # pic dimension
    page.showPDFpage(rect, imgPDF, 0)        # image fills the page
    psg.EasyProgressMeter("Import Images",   # show our progress
        i+1, imgcount)

doc.save("all-my-pics.pdf")

##Method 2: Embedding Files
#Script: PyMuPDF-Utilities-master.zip  all-my-pics-embedded.py
import os, fitz
import PySimpleGUI as psg                    # for showing progress bar
doc = fitz.open()                            # PDF with the pictures
imgdir = "D:/2012_10_05"                     # where the pictures are

imglist = os.listdir(imgdir)                 # list of pictures
imgcount = len(imglist)                      # pic count
imglist.sort()                               # nicely sort them

for i, f in enumerate(imglist):
    img = open(os.path.join(imgdir,f), "rb").read()    # make pic stream
    doc.embeddedFileAdd(img, f, filename=f,            # and embed it
                        ufilename=f, desc=f)
    psg.EasyProgressMeter("Embedding Files", # show our progress
        i+1, imgcount)

page = doc.newPage()                         # at least 1 page is needed,

doc.save("all-my-pics-embedded.pdf")

##Method 3: Attaching Files
#A third way to achieve this task is attaching files via page annotations 
#Script: PyMuPDF-Utilities-master.zip  all-my-pics-attached.py


##How to Interface with NumPy
#This shows how to create a PNG file from a numpy array 


import numpy as np
import fitz
#==============================================================================
# create a fun-colored width * height PNG with fitz and numpy
#==============================================================================
height = 150
width  = 100
bild = np.ndarray((height, width, 3), dtype=np.uint8)

for i in range(height):
    for j in range(width):
        # one pixel (some fun coloring)
        bild[i, j] = [(i+j)%256, i%256, j%256]

samples = bytearray(bild.tostring())    # get plain pixel data from numpy array
pix = fitz.Pixmap(fitz.csRGB, width, height, samples, alpha=False)
pix.writePNG("test.png")


##How to Extract Text in Natural Reading Order
#One of the common issues with PDF text extraction is, that text may not appear in any particular reading order.

#below script trys to overcome that 
from operator import itemgetter
from itertools import groupby
import fitz

def recover(words, rect):
    """ Word recovery.

    Notes:
        Method 'getTextWords()' does not try to recover words, if their single
        letters do not appear in correct lexical order. This function steps in
        here and creates a new list of recovered words.
    Args:
        words: list of words as created by 'getTextWords()'
        rect: rectangle to consider (usually the full page)
    Returns:
        List of recovered words. Same format as 'getTextWords', but left out
        block, line and word number - a list of items of the following format:
        [x0, y0, x1, y1, "word"]
    """
    # build my sublist of words contained in given rectangle
    mywords = [w for w in words if fitz.Rect(w[:4]) in rect]

    # sort the words by lower line, then by word start coordinate
    mywords.sort(key=itemgetter(3, 0))  # sort by y1, x0 of word rectangle

    # build word groups on same line
    grouped_lines = groupby(mywords, key=itemgetter(3))

    words_out = []  # we will return this

    # iterate through the grouped lines
    # for each line coordinate ("_"), the list of words is given
    for _, words_in_line in grouped_lines:
        for i, w in enumerate(words_in_line):
            if i == 0:  # store first word
                x0, y0, x1, y1, word = w[:5]
                continue

            r = fitz.Rect(w[:4])  # word rect

            # Compute word distance threshold as 20% of width of 1 letter.
            # So we should be safe joining text pieces into one word if they
            # have a distance shorter than that.
            threshold = r.width / len(w[4]) / 5
            if r.x0 <= x1 + threshold:  # join with previous word
                word += w[4]  # add string
                x1 = r.x1  # new end-of-word coordinate
                y0 = max(y0, r.y0)  # extend word rect upper bound
                continue

            # now have a new word, output previous one
            words_out.append([x0, y0, x1, y1, word])

            # store the new word
            x0, y0, x1, y1, word = w[:5]

        # output word waiting for completion
        words_out.append([x0, y0, x1, y1, word])

    return words_out

def search_for(text, words):
    """ Search for text in items of list of words

    Notes:
        Can be adjusted / extended in obvious ways, e.g. using regular
        expressions, or being case insensitive, or only looking for complete
        words, etc.
    Args:
        text: string to be searched for
        words: list of items in format delivered by 'getTextWords()'.
    Returns:
        List of rectangles, one for each found locations.
    """
    rect_list = []
    for w in words:
        if text in w[4]:
            rect_list.append(fitz.Rect(w[:4]))

    return rect_list


##How to Extract Tables from Documents
#If you see a table in a document, you are not normally looking 
#at something like an embedded Excel or other identifyable object. 
#It usually is just text, formatted to appear as appropriate.

#Script: PyMuPDF-Utilities-master.zip  wxTableExtract.py 


##How to Search for and Mark Text
#Instead of using  Page.searchFor(), using page.getTextWords

import sys
import fitz

def mark_word(page, text):
    """Underline each word that contains 'text'.
    """
    found = 0
    wlist = page.getTextWords()        # make the word list
    for w in wlist:                    # scan through all words on page
        if text in w[4]:               # w[4] is the word's string
            found += 1                 # count
            r = fitz.Rect(w[:4])       # make rect from word bbox
            page.addUnderlineAnnot(r)  # underline
    return found

fname = sys.argv[1]                    # filename
text = sys.argv[2]                     # search string
doc = fitz.open(fname)

print("underlining words containing '%s' in document '%s'" % (word, doc.name))

new_doc = False                        # indicator if anything found at all

for page in doc:                       # scan through the pages
    found = mark_word(page, text)      # mark the page's words
    if found:                          # if anything found ...
        new_doc = True
        print("found '%s' %i times on page %i" % (text, found, page.number + 1))

if new_doc:
    doc.save("marked-" + doc.name)

##How to Convert Any Document to PDF
#converts any PyMuPDF supported document to a PDF. 
#These include XPS, EPUB, FB2, CBZ and all image formats, including multi-page TIFF images.

from __future__ import print_function
"""
Demo script: Convert input file to a PDF
-----------------------------------------
Intended for multi-page input files like XPS, EPUB etc.

Features:
---------
Recovery of table of contents and links of input file.
While this works well for bookmarks (outlines, table of contents),
links will only work if they are not of type "LINK_NAMED".
This link type is skipped by the script.

For XPS and EPUB input, internal links however **are** of type "LINK_NAMED".
Base library MuPDF does not resolve them to page numbers.

So, for anyone expert enough to know the internal structure of these
document types, can further interpret and resolve these link types.

Dependencies
--------------
PyMuPDF v1.14.0+
"""
import sys
import fitz
if not (list(map(int, fitz.VersionBind.split("."))) >= [1,14,0]):
    raise SystemExit("need PyMuPDF v1.14.0+")
fn = sys.argv[1]

print("Converting '%s' to '%s.pdf'" % (fn, fn))

doc = fitz.open(fn)

b = doc.convertToPDF()                      # convert to pdf
pdf = fitz.open("pdf", b)                   # open as pdf

toc= doc.getToC()                           # table of contents of input
pdf.setToC(toc)                             # simply set it for output
meta = doc.metadata                         # read and set metadata
if not meta["producer"]:
    meta["producer"] = "PyMuPDF v" + fitz.VersionBind

if not meta["creator"]:
    meta["creator"] = "PyMuPDF PDF converter"
meta["modDate"] = fitz.getPDFnow()
meta["creationDate"] = meta["modDate"]
pdf.setMetadata(meta)

# now process the links
link_cnti = 0
link_skip = 0
for pinput in doc:                # iterate through input pages
    links = pinput.getLinks()     # get list of links
    link_cnti += len(links)       # count how many
    pout = pdf[pinput.number]     # read corresp. output page
    for l in links:               # iterate though the links
        if l["kind"] == fitz.LINK_NAMED:    # we do not handle named links
            print("named link page", pinput.number, l)
            link_skip += 1        # count them
            continue
        pout.insertLink(l)        # simply output the others

# save the conversion result
pdf.save(fn + ".pdf", garbage=4, deflate=True)
# say how many named links we skipped
if link_cnti > 0:
    print("Skipped %i named links of a total of %i in input." % (link_skip, link_cnti))

# now print any MuPDF warnings or errors:
errors = fitz.TOOLS.fitz_stderr
if errors:                        # any issues?
    print(errors)
    fitz.TOOLS.fitz_stderr_reset() # empty the message store

    
###PDF to Excel 
$ pip install tabula-py


#Example 
import tabula
import pandas as pd 

# Read pdf into DataFrame
df = tabula.read_pdf("test.pdf", pages='all')
df = tabula.read_pdf("data.pdf", pages="1-2,3")

list_of_df = tabula.read_pdf("test.pdf", pages='all', multiple_tables=True)
list_of_df = tabula.read_pdf("data.pdf", pages="1-2,3", multiple_tables=True)

# Read remote pdf into DataFrame
df2 = tabula.read_pdf("https://github.com/tabulapdf/tabula-java/raw/master/src/test/resources/technology/tabula/arabic.pdf")

# convert PDF into CSV
tabula.convert_into("test.pdf", "output.csv", output_format="csv", pages='all', multiple_tables=True)

# convert all PDFs in a directory
tabula.convert_into_by_batch("input_directory", output_format='csv', pages='all', multiple_tables=True)


##Options
    pages (str, int, list of int, optional)
        An optional values specifying pages to extract from. It allows str, int, list of int.
        Example: 1, '1-2,3', 'all' or [1,2]. Default is 1
    guess (bool, optional):
        Guess the portion of the page to analyze per page. Default True
    area (list of float, optional):
        Portion of the page to analyze(top,left,bottom,right).
        Example: [269.875, 12.75, 790.5, 561] or [[12.1,20.5,30.1,50.2],[1.0,3.2,10.5,40.2]]. Default is entire page
    relative_area (bool, optional):
        If all area values are between 0-100 (inclusive) and preceded by '%', input will be taken as % of actual height or width of the page. Default False.
    lattice (bool, optional):
        [spreadsheet option is deprecated] Force PDF to be extracted using lattice-mode extraction (if there are ruling lines separating each cell, as in a PDF of an Excel spreadsheet).
    stream (bool, optional):
        [nospreadsheet option is deprecated] Force PDF to be extracted using stream-mode extraction (if there are no ruling lines separating each cell, as in a PDF of an Excel spreadsheet)
    password (bool, optional):
        Password to decrypt document. Default is empty
    silent (bool, optional):
        Suppress all stderr output.
    columns (list, optional):
        X coordinates of column boundaries.
        Example: [10.1, 20.2, 30.3]
    output_format (str, optional):
        Format for output file or extracted object.
        For read_pdf(): json, dataframe
        For convert_into(): csv, tsv, json
    output_path (str, optional):
        Output file path. File format of it is depends on format.
        Same as --outfile option of tabula-java.
    java_options (list, optional):
        Set java options like -Xmx256m.
    pandas_options (dict, optional):
        Set pandas options like {'header': None}.
    multiple_tables (bool, optional):
        (Experimental) Extract multiple tables.
        This option uses JSON as an intermediate format, so if tabula-java output format will change, this option doesn't work.
    user_agent (str, optional)
        Set a custom user-agent when download a pdf from a url. 
        Otherwise it uses the default urllib.request user-agent
















